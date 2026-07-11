#!/usr/bin/env python3
"""Box-aware, fully-editable PPTX builder for the 3 architecture figures.

Unlike the SVG tracer, every card is ONE shape whose text frame wraps and is
vertically centered, so text never overflows and is evenly distributed. Cards
are spaced apart; small text 10pt, in-box titles 14pt; all text Microsoft YaHei
except math-script glyphs (𝒟 𝓘 𝓔 𝓛 𝒯), which keep their original look.

Usage: python3 pptx_build.py            # writes both agentic-runtime-{en,zh}.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.oxml.ns import qn

# 16:9 widescreen
SW, SH = Inches(13.333), Inches(7.5)
YAHEI = "Microsoft YaHei"

COL = {
  "skill":("E6F4EC","1A7D52","11603D"),
  "harn": ("ECE4FB","6D3FD4","4F2DA0"),
  "harnS":("F4EEFF","6D3FD4","4F2DA0"),
  "scaf": ("FBE7D3","BC5A16","9A4D12"),
  "data": ("DDEFF8","1474A6","0F5E87"),
  "neut": ("FFFDF8","CABFA8","221F1A"),
  "ok":   ("EAF7F0","1A7D52","11603D"),
}
INK=RGBColor(0x22,0x1f,0x1a); MUT=RGBColor(0x6b,0x63,0x57); PAGE=RGBColor(0xFF,0xFD,0xF8)
def C(name,i): return RGBColor.from_string(COL[name][i])

def _set_yahei(run):
    run.font.name = YAHEI
    rPr = run._r.get_or_add_rPr()
    for t in ('a:ea','a:cs'):
        e=rPr.find(qn(t))
        if e is None:
            e=rPr.makeelement(qn(t),{}); rPr.append(e)
        e.set('typeface',YAHEI)

def bg(slide):
    s=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,SW,SH)
    s.fill.solid(); s.fill.fore_color.rgb=PAGE; s.line.fill.background(); s.shadow.inherit=False
    return s

def heading(slide,title,sub):
    tb=slide.shapes.add_textbox(Inches(.45),Inches(.22),Inches(12.4),Inches(.9))
    tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; r=p.add_run(); r.text=title
    r.font.size=Pt(20); r.font.bold=True; r.font.color.rgb=INK; _set_yahei(r)
    p2=tf.add_paragraph(); r2=p2.add_run(); r2.text=sub
    r2.font.size=Pt(11); r2.font.color.rgb=MUT; _set_yahei(r2)

def card(slide,x,y,w,h,palette,*,title=None,title_size=14,lines=None,
         line_size=10,anchor=MSO_ANCHOR.MIDDLE,strong=False,fill=True,
         dash=False,bullets=False,title_color=None,line_color=None,sw=1.5,rxrel=None):
    fillc,strokec,textc=COL[palette]
    shp=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x),Inches(y),Inches(w),Inches(h))
    shp.shadow.inherit=False
    if fill:
        shp.fill.solid(); shp.fill.fore_color.rgb=RGBColor.from_string(COL["harnS"][0] if strong else fillc)
    else:
        shp.fill.background()
    shp.line.color.rgb=RGBColor.from_string(strokec); shp.line.width=Pt(sw)
    if dash:
        ln=shp.line._get_or_add_ln(); pd=ln.makeelement(qn('a:prstDash'),{'val':'dash'}); ln.append(pd)
    try: shp.adjustments[0]=0.08 if rxrel is None else rxrel
    except Exception: pass
    tf=shp.text_frame; tf.word_wrap=True
    try: tf.auto_size=MSO_AUTO_SIZE.NONE
    except Exception: pass
    tf.vertical_anchor=anchor
    for m in ('margin_left','margin_right'): setattr(tf,m,Inches(.14))
    for m in ('margin_top','margin_bottom'): setattr(tf,m,Inches(.08))
    tc=title_color or RGBColor.from_string(textc)
    lc=line_color or RGBColor.from_string(textc)
    first=True
    if title is not None:
        p=tf.paragraphs[0]; p.alignment=PP_ALIGN.LEFT
        r=p.add_run(); r.text=title; r.font.size=Pt(title_size); r.font.bold=True
        r.font.color.rgb=tc; _set_yahei(r); p.space_after=Pt(5); first=False
    for ln in (lines or []):
        p=tf.paragraphs[0] if first else tf.add_paragraph()
        first=False
        p.alignment=PP_ALIGN.LEFT; p.space_after=Pt(3); p.line_spacing=1.08
        r=p.add_run(); r.text=("•  "+ln if bullets else ln)
        r.font.size=Pt(line_size); r.font.color.rgb=lc; _set_yahei(r)
    return shp

def est_w_in(text,size_pt):
    """estimate rendered text width in inches"""
    w=0.0
    for c in text:
        w += size_pt*(1.0 if ord(c)>0x2E7F else 0.56)
    return w/72.0

def header_top(slide,bx,bw,top,title,size,color,*,badge_text=None,badge_col=None,badge_d=0.40):
    """Place title (and optional circled badge) horizontally CENTERED at the TOP
    of a box. Returns the y (inches) just below the header."""
    if badge_text is None:
        tb=slide.shapes.add_textbox(Inches(bx),Inches(top),Inches(bw),Inches(size/72*1.6))
        tf=tb.text_frame; tf.word_wrap=False
        for m in ('margin_left','margin_right','margin_top','margin_bottom'): setattr(tf,m,0)
        p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
        r=p.add_run(); r.text=title; r.font.size=Pt(size); r.font.bold=True
        r.font.color.rgb=color; _set_yahei(r)
        return top+size/72*1.7
    tw=est_w_in(title,size); sp=0.12; gw=badge_d+sp+tw
    sx=bx+(bw-gw)/2
    badge(slide,sx,top,badge_d,badge_text,badge_col)
    tb=slide.shapes.add_textbox(Inches(sx+badge_d+sp),Inches(top-0.03),Inches(tw+0.5),Inches(badge_d+0.06))
    tf=tb.text_frame; tf.word_wrap=False; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    for m in ('margin_left','margin_right','margin_top','margin_bottom'): setattr(tf,m,0)
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.LEFT
    r=p.add_run(); r.text=title; r.font.size=Pt(size); r.font.bold=True
    r.font.color.rgb=color; _set_yahei(r)
    return top+badge_d

def body_box(slide,x,y,w,h,items,*,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP):
    """items: list of (text,size,color,bold). Wraps, never overflows."""
    tb=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf=tb.text_frame; tf.word_wrap=True
    try: tf.auto_size=MSO_AUTO_SIZE.NONE
    except Exception: pass
    tf.vertical_anchor=anchor
    tf.margin_left=Inches(.04); tf.margin_right=Inches(.04); tf.margin_top=0; tf.margin_bottom=0
    first=True
    for text,size,color,bold in items:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        p.alignment=align; p.line_spacing=1.06; p.space_after=Pt(3)
        r=p.add_run(); r.text=text; r.font.size=Pt(size); r.font.bold=bold
        r.font.color.rgb=color; _set_yahei(r)
    return tb

def badge(slide,x,y,d,text,colhex):
    s=slide.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x),Inches(y),Inches(d),Inches(d))
    s.shadow.inherit=False; s.fill.solid(); s.fill.fore_color.rgb=RGBColor.from_string(colhex)
    s.line.color.rgb=RGBColor(0xff,0xff,0xff); s.line.width=Pt(1.5)
    tf=s.text_frame; tf.word_wrap=False; p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text=text; r.font.size=Pt(11); r.font.bold=True
    r.font.color.rgb=RGBColor(0xff,0xff,0xff); _set_yahei(r)
    return s

def arrow(slide,x1,y1,x2,y2,colhex,*,dash=False,elbow=False,sw=1.8):
    kind=MSO_CONNECTOR.ELBOW if elbow else MSO_CONNECTOR.STRAIGHT
    cn=slide.shapes.add_connector(kind,Inches(x1),Inches(y1),Inches(x2),Inches(y2))
    cn.line.color.rgb=RGBColor.from_string(colhex); cn.line.width=Pt(sw)
    ln=cn.line._get_or_add_ln()
    te=ln.makeelement(qn('a:tailEnd'),{'type':'triangle','w':'med','len':'med'}); ln.append(te)
    if dash:
        pd=ln.makeelement(qn('a:prstDash'),{'val':'dash'}); ln.append(pd)
    return cn

# =================================================================
# CONTENT
# =================================================================
TXT={
"en":{
 "f1":{
   "title":"The Three Layers — What Each One Is For",
   "sub":"A first look: an AI agent is built from three stacked layers, plus a data subsystem on the side.",
   "bands":[
     ("skill","L3","Skills  —  “What to do”",[
        "Holds the job descriptions for tasks: the goal, the input / output format, and the rules for when to act, stop and loop.",
        "Like an employee’s playbook / SOP.",
        "Add a Skill → the agent can do one more kind of task  (logical scaling)"]),
     ("harn","L2","Harness  —  “How the work gets done”",[
        "Turns a Skill’s intent into concrete, runnable steps. Manages memory, context, tools, code generation, safety checks and audit logs.",
        "Like a project manager + the runtime framework.",
        "The single place where Skills and Scaffold meet  (the decoupling point)"]),
     ("scaf","L1","Scaffold  —  “Where it runs, safely”",[
        "Provides the isolated place to execute: sandbox, shell, network, storage and model-serving capacity, with security boundaries.",
        "Like a secure cloud computer.",
        "Add an instance → more throughput  (physical scaling)"]),
   ],
   "data_title":"Data Subsystem 𝒟",
   "data_sub":"What the agent knows about your data (kept off to the side, mostly offline)",
   "data":["D₁  fetch — read the slice you need, on demand",
           "D₂  semantic summary — pre-digest sources offline",
           "D₃  governance — remember how to use the data",
           "D₄  lifetime — freshness, scope & cost budget",
           "Ω   workspace — run history, artifacts, LLM wiki"],
   "foot":"Key idea: each kind of growth (more skills / more machines / more data) is added on a different axis, so they don’t slow each other down.",
 },
 "f2":{
   "title":"What Each Layer Actually Contains — the Spec Lists",
   "sub":"The concrete, buildable pieces in every layer, grounded in real formats: Anthropic SKILL.md, MCP tools, Claude Code, OpenAI connectors.",
   "l3":("skill","L3 · Skills","SkillSpec — Anthropic SKILL.md format",
        ["name + description (when to use)","input / output schema",
         "output format + constraints","call / stop / loop conditions",
         "examples (shots) + anti-examples","personalization + verification"]),
   "l2":("harn","L2 · Harness","CapabilityCapsule — MCP-grounded",
        ["CapabilityCapsule: inputSchema · contract · invariants",
         "tool routing (intent-scoped, top-k)",
         "context assembly (shots · schema · memory)",
         "code generation (Read / Edit / Bash + test / lint)",
         "tool synthesis (sandbox-validated capsule)",
         "memory maintenance M (compaction · writeback · reflection · retrieval · GC)",
         "reward / verifier + policy gate + audit log"]),
   "l1":("scaf","L1 · Scaffold","ExecutionSpec — isolation & resources",
        ["isolation: microVM · container · wasm",
         "bash / process / filesystem",
         "network: allowlist · egress policy",
         "SSO / identity: OIDC · SAML · SCIM · RBAC",
         "cloud: AWS / GCP / Azure · KMS · S3 · PrivateLink",
         "serving (KV cache · batching) · lifecycle · observability"]),
   "data_title":"Data Subsystem 𝒟",
   "data":["D₁ DataSourceCard — source · auth · quota · NFR",
           "D₂ SemanticSummary — schema* · relations · themes",
           "D₃ DataUsageSkill — sources · join plan · runs",
           "D₄ Lifetime — freshness · retention · budget",
           "Ω  Workspace — RunTrace · ArtifactManifest · LLM Wiki"],
 },
 "f3":{
   "title":"How the Agent Works — a Step-by-Step Walkthrough",
   "sub":"The Harness runs this loop. The LLM makes the choices; typed contracts, a policy gate and a verifier keep it safe and on track.",
   "steps":[
     ("1","Understand the goal","Read the request and constraints; figure out what “done” looks like."),
     ("2","Pick the skill","Choose the SkillSpec that fits and load its input / output schema."),
     ("3","Gather context","Pull in memory, relevant examples (shots) and the tools’ contracts."),
     ("4","Make a plan","The LLM drafts a plan with fallbacks and a clear stop condition."),
     ("5","Choose the next action","Pick the tool (O1), the model (O7) and the token budget (O8)."),
     ("6","Check, then run","Policy gate checks scope & permissions, then run on Scaffold / query Data."),
     ("7","Observe & score","Collect result and evidence; a verifier / reward grades it."),
     ("8","Good enough?","Yes → deliver the result.   No → reflect, revise, loop back to step 5."),
   ],
   "scaf_title":"Scaffold",
   "scaf":["sandbox · shell · serving","filesystem · network · SSO / cloud"],
   "data_title":"Data subsystem 𝒟",
   "data":["D₁ / D₂ query · D₃ / D₄ policy","Ω workspace evidence"],
   "deliver":"Deliver artifact",
   "loopback":"reflect & loop back to step 5",
   "sysnote":"Running quietly in the background: System skills (M) compact, write back and retrieve memory between steps.",
 },
},
"zh":{
 "f1":{
   "title":"三个层 — 每一层是干什么的",
   "sub":"第一眼认识：一个 AI agent 由三个层叠起来，旁边再加一个数据子系统。",
   "bands":[
     ("skill","L3","Skills  —  “要做什么”",[
        "存放任务的“岗位说明书”：目标、输入 / 输出格式，以及何时行动、何时停止、何时循环的规则。",
        "就像员工的工作手册 / SOP。",
        "加一个 Skill → 多一种能做的任务（逻辑扩展）"]),
     ("harn","L2","Harness  —  “怎么把活干成”",[
        "把 Skill 的意图翻译成可执行的步骤。管理记忆、上下文、工具、代码生成、安全检查与审计日志。",
        "就像项目经理 + 运行时框架。",
        "Skills 与 Scaffold 唯一相遇的地方（解耦点）"]),
     ("scaf","L1","Scaffold  —  “在哪里安全地跑”",[
        "提供隔离的执行环境：sandbox、shell、网络、存储与模型 serving 算力，并带安全边界。",
        "就像一台安全的云电脑。",
        "加一个实例 → 多一份吞吐（物理扩展）"]),
   ],
   "data_title":"数据子系统 𝒟",
   "data_sub":"agent 对你数据的认知（放在一旁，主要离线）",
   "data":["D₁  取数 — 按需只读你要的那一片",
           "D₂  语义总结 — 离线把数据源预消化",
           "D₃  治理记忆 — 记住“该怎么用这份数据”",
           "D₄  lifetime — 时新性、口径与成本预算",
           "Ω   工作区 — 运行历史、产物、LLM wiki"],
   "foot":"关键点：每一种增长（更多 skill / 更多机器 / 更多数据）加在不同的轴上，因此互不拖累。",
 },
 "f2":{
   "title":"每一层里到底装了什么 — Spec 清单",
   "sub":"每层可落地、可构建的具体部件，对齐真实格式：Anthropic SKILL.md、MCP tools、Claude Code、OpenAI connectors。",
   "l3":("skill","L3 · Skills","SkillSpec — Anthropic SKILL.md 格式",
        ["name + description（何时使用）","input / output schema",
         "output format + 约束","call / stop / loop 条件",
         "examples（shots）+ 反例","个性化 + verification"]),
   "l2":("harn","L2 · Harness","CapabilityCapsule — 以 MCP 为底座",
        ["CapabilityCapsule：inputSchema · contract · invariants",
         "工具路由（intent-scoped, top-k）",
         "context 组装（shots · schema · memory）",
         "code generation（Read / Edit / Bash + test / lint）",
         "tool synthesis（沙箱验证的 capsule）",
         "记忆维护 M（compaction · writeback · reflection · retrieval · GC）",
         "reward / verifier + policy gate + 审计日志"]),
   "l1":("scaf","L1 · Scaffold","ExecutionSpec — 隔离与资源",
        ["隔离：microVM · container · wasm",
         "bash / process / filesystem",
         "网络：allowlist · egress 策略",
         "SSO / 身份：OIDC · SAML · SCIM · RBAC",
         "云：AWS / GCP / Azure · KMS · S3 · PrivateLink",
         "serving（KV cache · batching）· lifecycle · 观测"]),
   "data_title":"数据子系统 𝒟",
   "data":["D₁ DataSourceCard — source · auth · quota · NFR",
           "D₂ SemanticSummary — schema* · relations · themes",
           "D₃ DataUsageSkill — sources · join plan · runs",
           "D₄ Lifetime — freshness · retention · budget",
           "Ω  工作区 — RunTrace · ArtifactManifest · LLM Wiki"],
 },
 "f3":{
   "title":"Agent 是怎么工作的 — 分步走查",
   "sub":"Harness 运行这个循环。由 LLM 做选择；typed contract、policy gate 与 verifier 保证它安全、不跑偏。",
   "steps":[
     ("1","读懂目标","读取请求与约束；想清楚“做完”长什么样。"),
     ("2","挑选 skill","选出合适的 SkillSpec，加载它的输入 / 输出 schema。"),
     ("3","收集上下文","拉入记忆、相关示例（shots）和工具的契约。"),
     ("4","制定计划","LLM 起草带 fallback 和明确停止条件的计划。"),
     ("5","选择下一步动作","为这一步选 tool（O1）、模型（O7）和 token 预算（O8）。"),
     ("6","先检查，再执行","policy gate 校验范围与权限，然后在 Scaffold 执行 / 查询 Data。"),
     ("7","观察并打分","收集结果与证据；verifier / reward 给它打分。"),
     ("8","够好了吗？","够好 → 交付结果。   不够 → 反思、修订，回到第 5 步。"),
   ],
   "scaf_title":"Scaffold",
   "scaf":["sandbox · shell · serving","filesystem · network · SSO / cloud"],
   "data_title":"数据子系统 𝒟",
   "data":["D₁ / D₂ 查询 · D₃ / D₄ 策略","Ω workspace 证据"],
   "deliver":"交付产物",
   "loopback":"反思并回到第 5 步",
   "sysnote":"在后台静默运行：System skills (M) 在步骤之间做记忆压缩、写回与检索。",
 },
},
}

# =================================================================
# LAYOUTS
# =================================================================
def slide_fig1(prs,t):
    s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s); heading(s,t["title"],t["sub"])
    LX,LW=0.45,7.55
    bandH,gap,y0=1.62,0.22,1.35
    cols=["skill","harn","scaf"]
    ys=[y0+i*(bandH+gap) for i in range(3)]
    for i,(pal,tag,name,body) in enumerate(t["bands"]):
        y=ys[i]
        card(s,LX,y,LW,bandH,pal,title=None,lines=None)  # base box
        # header: badge + title centered at top
        by=header_top(s,LX,LW,y+0.12,name,14,C(pal,2),badge_text=tag,badge_col=COL[pal][1])
        # body below header, centered block
        items=[]
        for j,ln in enumerate(body):
            if j>=len(body)-2: items.append((ln,10,C(pal,1),True))
            else: items.append((ln,10,C(pal,2),False))
        body_box(s,LX+0.3,by+0.06,LW-0.6,bandH-(by-y)-0.16,items,anchor=MSO_ANCHOR.MIDDLE)
    # data column
    DX,DW=8.3,4.55; DH=ys[2]+bandH-y0
    card(s,DX,y0,DW,DH,"data")  # base box only
    dby=header_top(s,DX,DW,y0+0.16,t["data_title"],14,C("data",2))
    ditems=[(t["data_sub"],10,C("data",2),False)]+[(x,10,C("data",2),False) for x in t["data"]]
    body_box(s,DX+0.28,dby+0.04,DW-0.56,DH-(dby-y0)-0.2,ditems,anchor=MSO_ANCHOR.MIDDLE)
    # arrows between bands (right lane)
    ax=LX+LW-0.35
    arrow(s,ax,ys[0]+bandH,ax,ys[1],COL["harn"][1])
    arrow(s,ax,ys[1]+bandH,ax,ys[2],COL["scaf"][1])
    # harness <-> data
    arrow(s,LX+LW,ys[1]+bandH*0.5,DX,ys[1]+bandH*0.5,COL["data"][1])
    # footer
    fy=ys[2]+bandH+0.16
    card(s,LX,fy,DX+DW-LX,7.45-fy,"neut",lines=[t["foot"]],line_size=11,
         anchor=MSO_ANCHOR.MIDDLE,dash=True)

def slide_fig2(prs,t):
    s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s); heading(s,t["title"],t["sub"])
    LX,LW=0.45,7.55; y0=1.32
    specs=[t["l3"],t["l2"],t["l1"]]
    heights=[1.66,2.15,1.72]; gap=0.2
    ys=[y0];
    for hh in heights[:-1]: ys.append(ys[-1]+hh+gap)
    for (pal,name,note,items),y,hh in zip(specs,ys,heights):
        card(s,LX,y,LW,hh,pal)
        by=header_top(s,LX,LW,y+0.13,name+"   —   "+note,14,C(pal,2))
        bitems=[("•  "+it,10,C(pal,2),False) for it in items]
        body_box(s,LX+0.34,by+0.04,LW-0.6,hh-(by-y)-0.14,bitems,anchor=MSO_ANCHOR.MIDDLE)
    # arrows
    ax=LX+LW-0.3
    arrow(s,ax,ys[0]+heights[0],ax,ys[1],COL["harn"][1])
    arrow(s,ax,ys[1]+heights[1],ax,ys[2],COL["scaf"][1])
    # data
    DX,DW=8.3,4.55; DH=ys[2]+heights[2]-y0
    card(s,DX,y0,DW,DH,"data")
    dby=header_top(s,DX,DW,y0+0.16,t["data_title"],14,C("data",2))
    ditems=[("•  "+it,10,C("data",2),False) for it in t["data"]]
    body_box(s,DX+0.3,dby+0.06,DW-0.56,DH-(dby-y0)-0.2,ditems,anchor=MSO_ANCHOR.MIDDLE)
    arrow(s,LX+LW,ys[1]+heights[1]*0.5,DX,ys[1]+heights[1]*0.5,COL["data"][1])

def slide_fig3(prs,t):
    s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s); heading(s,t["title"],t["sub"])
    SX,SW=2.45,5.45; y0=1.22; stepH=0.70; gap=0.075
    ys=[y0+i*(stepH+gap) for i in range(8)]
    palmap=["skill","skill","data","harn","harn","harn","harn","scaf"]
    for i,(num,tt,desc) in enumerate(t["steps"]):
        y=ys[i]; pal=palmap[i]; strong=(i in (3,4,7))
        card(s,SX,y,SW,stepH,"harn",strong=strong,sw=1.4)
        # badge + title centered at TOP of the step box
        by=header_top(s,SX,SW,y+0.07,tt,13,C("harn",2),badge_text=num,badge_col=COL[pal][1],badge_d=0.34)
        # description below, centered horizontally
        body_box(s,SX+0.3,by+0.0,SW-0.6,stepH-(by-y)-0.06,
                 [(desc,10,MUT,False)],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.TOP)
    for i in range(7):
        arrow(s,SX+SW/2,ys[i]+stepH,SX+SW/2,ys[i+1],COL["harn"][1],sw=1.5)
    # right column: scaffold + data aligned around steps 5-7
    RX,RW=8.25,4.6
    sc_y=ys[4]
    card(s,RX,sc_y,RW,1.12,"scaf")
    scby=header_top(s,RX,RW,sc_y+0.1,t["scaf_title"],14,C("scaf",2))
    body_box(s,RX+0.28,scby+0.02,RW-0.56,1.12-(scby-sc_y)-0.14,
             [(x,10,C("scaf",2),False) for x in t["scaf"]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    da_y=sc_y+1.28
    card(s,RX,da_y,RW,1.12,"data")
    daby=header_top(s,RX,RW,da_y+0.1,t["data_title"],14,C("data",2))
    body_box(s,RX+0.28,daby+0.02,RW-0.56,1.12-(daby-da_y)-0.14,
             [(x,10,C("data",2),False) for x in t["data"]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    # step6 -> scaffold/data
    arrow(s,SX+SW,ys[5]+stepH/2,RX,sc_y+0.56,COL["scaf"][1],elbow=True)
    arrow(s,SX+SW,ys[6]+stepH/2,RX,da_y+0.56,COL["data"][1],elbow=True)
    # deliver box (bottom right), step8 -> deliver
    dl_y=ys[7]
    card(s,RX,dl_y,RW,stepH,"ok",title=None,lines=[t["deliver"]],line_size=14,
         anchor=MSO_ANCHOR.MIDDLE)
    arrow(s,SX+SW,ys[7]+stepH/2,RX,dl_y+stepH/2,COL["skill"][1])
    # loopback far-left lane: step8 -> step5
    lane=2.0
    arrow(s,SX,ys[7]+stepH/2,lane,ys[7]+stepH/2,COL["scaf"][1],dash=True)
    arrow(s,lane,ys[7]+stepH/2,lane,ys[4]+stepH/2,COL["scaf"][1],dash=True)
    arrow(s,lane,ys[4]+stepH/2,SX,ys[4]+stepH/2,COL["scaf"][1],dash=True)
    tb=s.shapes.add_textbox(Inches(0.55),Inches((ys[4]+ys[7])/2),Inches(1.4),Inches(0.5))
    tf=tb.text_frame; tf.word_wrap=True; p=tf.paragraphs[0]; r=p.add_run(); r.text=t["loopback"]
    r.font.size=Pt(10); r.font.bold=True; r.font.color.rgb=C("scaf",1); _set_yahei(r)
    # system-skills note band
    ny=ys[7]+stepH+0.12
    card(s,SX,ny,RX+RW-SX,min(0.5,7.34-ny),"neut",lines=[t["sysnote"]],line_size=10,
         anchor=MSO_ANCHOR.MIDDLE,dash=True)

def build(lang):
    prs=Presentation(); prs.slide_width=SW; prs.slide_height=SH
    slide_fig1(prs,TXT[lang]["f1"])
    slide_fig2(prs,TXT[lang]["f2"])
    slide_fig3(prs,TXT[lang]["f3"])
    out=f"agentic-runtime-{lang}.pptx"; prs.save(out); print("wrote",out)

if __name__=="__main__":
    for lg in ("en","zh"): build(lg)
