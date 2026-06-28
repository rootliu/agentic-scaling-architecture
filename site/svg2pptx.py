#!/usr/bin/env python3
"""Convert the generated figure SVGs into native, editable PowerPoint shapes.

Each <rect>/<text>/<path>/<circle> becomes a real PPTX object (rounded rectangle,
text box, connector/freeform line with arrowhead, oval) so everything is movable,
recolorable and re-typable inside PowerPoint. One figure per slide.

Usage:  python3 svg2pptx.py <lang>      # lang = en | zh
Outputs <lang> deck next to this script as agentic-runtime-<lang>.pptx
"""
import sys, re, math, xml.etree.ElementTree as ET
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.oxml.ns import qn

VBW, VBH = 1200, 860            # svg viewBox
# slide sized to the figure aspect, 12in wide
SLIDE_W_IN = 12.0
PX = Emu(int(914400 * SLIDE_W_IN / VBW))   # EMU per svg px
def E(px):  # px -> EMU
    return Emu(int(round(px * 914400 * SLIDE_W_IN / VBW)))
SLIDE_W = E(VBW)
SLIDE_H = E(VBH)

NS = "{http://www.w3.org/2000/svg}"

def rgb(s):
    if not s or s == "none":
        return None
    s = s.strip()
    if s.startswith("#"):
        s = s[1:]
        if len(s) == 3:
            s = "".join(c*2 for c in s)
        return RGBColor(int(s[0:2],16), int(s[2:4],16), int(s[4:6],16))
    return None

def add_arrowhead(line_elem, color, end=True):
    """Add a filled triangular arrowhead to a line/freeform via raw XML."""
    ln = line_elem
    tag = "tailEnd" if end else "headEnd"
    e = ln.makeelement(qn('a:'+tag), {'type':'triangle','w':'med','len':'med'})
    ln.append(e)

def parse_path_points(d):
    """Return list of (x,y) for a path made only of M/L commands."""
    toks = re.findall(r'[ML]\s*(-?\d+(?:\.\d+)?)\s+(-?\d+(?:\.\d+)?)', d)
    return [(float(a), float(b)) for a,b in toks]

def parse_path_cubic_endpoints(d):
    """For C (cubic) paths we only need start + final endpoint for a straight-ish
    connector; but to keep elbows we approximate by sampling all coordinate pairs."""
    nums = re.findall(r'(-?\d+(?:\.\d+)?)', d)
    pts = [(float(nums[i]), float(nums[i+1])) for i in range(0, len(nums)-1, 2)]
    return pts

def build(lang):
    tree = ET.parse(f"figures/scene-base.svg")  # placeholder, replaced per-fig below
    raise SystemExit

def make_deck(lang):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]
    for fig in ["f1","f2","f3"]:
        svg = open(f"figures/{fig}-{lang}.svg", encoding="utf-8").read()
        # strip namespace decl for simpler tag handling
        root = ET.fromstring(svg)
        slide = prs.slides.add_slide(blank)
        render_svg(root, slide)
    out = f"agentic-runtime-{lang}.pptx"
    prs.save(out)
    print("wrote", out)

def iter_all(root):
    for el in root.iter():
        yield el

def render_svg(root, slide):
    shapes = slide.shapes
    # background white
    bg = shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(0xFF,0xFD,0xF8)
    bg.line.fill.background()
    bg.shadow.inherit = False

    for el in iter_all(root):
        tag = el.tag.replace(NS,"")
        if tag == "rect":
            x=float(el.get("x",0)); y=float(el.get("y",0))
            w=float(el.get("width",0)); h=float(el.get("height",0))
            if w>=VBW-2 and h>=VBH-2:   # skip the full-canvas bg/grid rects
                continue
            rx=float(el.get("rx",0) or 0)
            shape = shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE if rx>0 else MSO_SHAPE.RECTANGLE,
                E(x),E(y),E(w),E(h))
            shape.shadow.inherit=False
            f=rgb(el.get("fill"))
            if f is None or el.get("fill")=="none":
                shape.fill.background()
            else:
                shape.fill.solid(); shape.fill.fore_color.rgb=f
            st=rgb(el.get("stroke"))
            if st is None:
                shape.line.fill.background()
            else:
                shape.line.color.rgb=st
                sw=float(el.get("stroke-width",1) or 1)
                shape.line.width=Pt(sw*0.75)
                if el.get("stroke-dasharray"):
                    d=shape.line._get_or_add_ln()
                    pd=d.makeelement(qn('a:prstDash'),{'val':'dash'}); d.append(pd)
            # rounded corner radius
            if rx>0:
                try:
                    shape.adjustments[0]=min(0.5, rx/min(w,h))
                except Exception: pass
        elif tag == "circle":
            cx=float(el.get("cx",0)); cy=float(el.get("cy",0)); r=float(el.get("r",0))
            shape=shapes.add_shape(MSO_SHAPE.OVAL,E(cx-r),E(cy-r),E(2*r),E(2*r))
            shape.shadow.inherit=False
            f=rgb(el.get("fill"))
            if f: shape.fill.solid(); shape.fill.fore_color.rgb=f
            else: shape.fill.background()
            st=rgb(el.get("stroke"))
            if st: shape.line.color.rgb=st; shape.line.width=Pt(float(el.get("stroke-width",1))*0.75)
            else: shape.line.fill.background()
        elif tag == "path":
            d=el.get("d","")
            stroke=rgb(el.get("stroke"))
            if stroke is None: continue
            if "C" in d:
                pts=parse_path_cubic_endpoints(d)
                pts=[pts[0],pts[-1]]   # simplify cubic to straight connector
            else:
                pts=parse_path_points(d)
            if len(pts)<2: continue
            draw_polyline(shapes,pts,stroke,float(el.get("stroke-width",2) or 2),
                          bool(el.get("marker-end")), bool(el.get("stroke-dasharray")))
        elif tag == "text":
            add_text(shapes, el)

def draw_polyline(shapes, pts, color, sw, arrow, dashed):
    # use a freeform through the points
    fb = shapes.build_freeform(E(pts[0][0]), E(pts[0][1]), scale=1.0)
    fb.add_line_segments([(E(x),E(y)) for x,y in pts[1:]], close=False)
    shp = fb.convert_to_shape()
    shp.fill.background()
    shp.line.color.rgb=color
    shp.line.width=Pt(sw*0.75)
    ln=shp.line._get_or_add_ln()
    if dashed:
        pd=ln.makeelement(qn('a:prstDash'),{'val':'dash'}); ln.append(pd)
    if arrow:
        add_arrowhead(ln,color,end=True)
    shp.shadow.inherit=False

def add_text(shapes, el):
    txt=("".join(el.itertext())).strip()
    if not txt: return
    x=float(el.get("x",0)); y=float(el.get("y",0))
    size=float(el.get("font-size",13) or 13)
    anchor=el.get("text-anchor","start")
    baseline=el.get("dominant-baseline","")
    weight=el.get("font-weight","400")
    fill=rgb(el.get("fill")) or RGBColor(0x22,0x1f,0x1a)
    # estimate width (generous so single-line text never wraps)
    cjk = any('一'<=c<='鿿' for c in txt)
    cw = size*(1.12 if cjk else 0.62)
    w = cw*len(txt) + size*1.2 + 8
    h = size*1.6
    if anchor=="middle": left=x-w/2
    elif anchor=="end":  left=x-w
    else:                left=x-4
    # vertical: svg y is baseline (or middle); textbox is top-anchored
    if baseline=="middle": top=y-h/2
    else:                  top=y-size*0.98
    tb=shapes.add_textbox(E(left),E(top),E(w),E(h))
    tf=tb.text_frame; tf.word_wrap=False
    try: tf.auto_size=MSO_AUTO_SIZE.NONE
    except Exception: pass
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    # prevent PowerPoint from shrinking/wrapping: keep frame on one line
    bodyPr=tf._txBody.find(qn('a:bodyPr'))
    if bodyPr is not None:
        bodyPr.set('wrap','none')
    p=tf.paragraphs[0]
    p.alignment={"middle":PP_ALIGN.CENTER,"end":PP_ALIGN.RIGHT}.get(anchor,PP_ALIGN.LEFT)
    r=p.add_run(); r.text=txt
    fnt=r.font
    # two fixed tiers: large in-box text 14pt, small embedded text 12pt
    is_large = size >= 15
    fnt.size = Pt(14 if is_large else 12)
    fnt.bold = int(re.sub(r"\D","",weight) or 400) >= 600
    fnt.color.rgb=fill
    # all non-math text -> Microsoft YaHei; math-script glyphs (𝒟 𝓘 𝓔 𝓛 𝒯 Θ …)
    # fall back automatically and keep their original 花体 appearance.
    YAHEI = "Microsoft YaHei"
    fnt.name = YAHEI
    rPr = r._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', YAHEI)

if __name__=="__main__":
    lang=sys.argv[1] if len(sys.argv)>1 else "zh"
    make_deck(lang)
