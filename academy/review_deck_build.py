#!/usr/bin/env python3
"""架构评审 deck 构建器 — v25 合回单篇方案。

产物: academy/architecture-review-v25-zh.pptx (16:9, 中文)
用法: python3 academy/review_deck_build.py

配色沿用 site/style.css 迁移后的 Tailwind 系:
  logic #3B82F6 / contract #7C3AED / phys #EA580C / data #059669
每张卡片是一个独立 shape，文本框自动换行并垂直居中，便于评审现场直接编辑。
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

SW, SH = Inches(13.333), Inches(7.5)
YAHEI = "Microsoft YaHei"

# (底色, 描边, 文字)
PAL = {
    "logic":    ("EFF6FF", "3B82F6", "1E40AF"),
    "contract": ("F5F3FF", "7C3AED", "5B21B6"),
    "phys":     ("FFF7ED", "EA580C", "C2410C"),
    "data":     ("ECFDF5", "059669", "047857"),
    "neut":     ("F9FAFB", "E5E7EB", "374151"),
    "warn":     ("FEF2F2", "DC2626", "991B1B"),
    "ink":      ("FFFFFF", "D1D5DB", "111827"),
}
INK = RGBColor(0x11, 0x18, 0x27)
MUT = RGBColor(0x6B, 0x72, 0x80)
PAGE = RGBColor(0xFF, 0xFF, 0xFF)


def C(name, i):
    return RGBColor.from_string(PAL[name][i])


def _yahei(run):
    run.font.name = YAHEI
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {})
            rPr.append(e)
        e.set("typeface", YAHEI)


def blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bgfill = s.background.fill
    bgfill.solid()
    bgfill.fore_color.rgb = PAGE
    return s


def textbox(slide, x, y, w, h, lines, *, size=11, color=INK, bold=False,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        p.space_after = Pt(3)
        txt, sz, bd, col = ln if isinstance(ln, tuple) else (ln, size, bold, color)
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(sz)
        r.font.bold = bd
        r.font.color.rgb = col
        _yahei(r)
    return tb


def heading(slide, title, sub=None, *, tag=None, tagpal="contract"):
    textbox(slide, Inches(0.62), Inches(0.40), Inches(10.6), Inches(0.62),
            [(title, 27, True, INK)])
    if sub:
        textbox(slide, Inches(0.62), Inches(1.02), Inches(11.4), Inches(0.42),
                [(sub, 12, False, MUT)])
    if tag:
        w = Inches(0.30 + 0.115 * len(tag))
        s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   SW - Inches(0.62) - w, Inches(0.44), w, Inches(0.34))
        s.fill.solid()
        s.fill.fore_color.rgb = C(tagpal, 0)
        s.line.color.rgb = C(tagpal, 1)
        s.line.width = Pt(1.0)
        try:
            s.adjustments[0] = 0.22
        except Exception:
            pass
        tf = s.text_frame
        tf.word_wrap = False
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = tag
        r.font.size = Pt(10)
        r.font.bold = True
        r.font.color.rgb = C(tagpal, 2)
        _yahei(r)
    # 细分隔线
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.62), Inches(1.50),
                                SW - Inches(1.24), Pt(1.1))
    ln.fill.solid()
    ln.fill.fore_color.rgb = RGBColor.from_string("E5E7EB")
    ln.line.fill.background()
    ln.shadow.inherit = False


def card(slide, x, y, w, h, pal, *, title=None, lines=None, title_size=14,
         body_size=10.5, dash=False):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = C(pal, 0)
    s.line.color.rgb = C(pal, 1)
    s.line.width = Pt(1.2)
    s.shadow.inherit = False
    if dash:
        lnEl = s.line._get_or_add_ln()
        pd = lnEl.makeelement(qn("a:prstDash"), {"val": "dash"})
        lnEl.append(pd)
    try:
        s.adjustments[0] = 0.06
    except Exception:
        pass
    tf = s.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = Inches(0.13)
    tf.margin_top = Inches(0.10)
    tf.margin_bottom = Inches(0.08)
    first = True
    if title:
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = title
        r.font.size = Pt(title_size)
        r.font.bold = True
        r.font.color.rgb = C(pal, 2)
        _yahei(r)
        first = False
    for ln in (lines or []):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.14
        p.space_before = Pt(3)
        txt, sz = (ln, body_size) if isinstance(ln, str) else ln
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(sz)
        r.font.color.rgb = INK
        _yahei(r)
    return s


def footer(slide, n, total, note=""):
    textbox(slide, Inches(0.62), SH - Inches(0.50), Inches(9.0), Inches(0.30),
            [(note, 9, False, MUT)])
    textbox(slide, SW - Inches(1.55), SH - Inches(0.50), Inches(0.95), Inches(0.30),
            [(f"{n} / {total}", 9, False, MUT)], align=PP_ALIGN.RIGHT)


# ─────────────────────────────── slides ───────────────────────────────

TOTAL = 16


def s01_cover(prs):
    s = blank(prs)
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.22), SH)
    band.fill.solid()
    band.fill.fore_color.rgb = C("contract", 1)
    band.line.fill.background()
    band.shadow.inherit = False
    textbox(s, Inches(0.95), Inches(2.10), Inches(11.2), Inches(0.5),
            [("架构评审 · v25", 15, True, C("contract", 2))])
    textbox(s, Inches(0.95), Inches(2.62), Inches(11.4), Inches(1.5),
            [("契约中心的可扩展 · 可管理", 40, True, INK),
             ("Agentic Runtime 参考架构", 40, True, INK)])
    textbox(s, Inches(0.95), Inches(4.34), Inches(11.0), Inches(0.9),
            [("A Contract-Centered Architecture for Scalable and Manageable Agentic Runtimes", 13, False, MUT),
             ("本次评审议题：撤销 v21/v22 拆分，合回单篇 v25", 13, False, C("phys", 2))])
    card(s, Inches(0.95), Inches(5.55), Inches(3.42), Inches(1.05), "logic",
         title="责任模型", lines=["Skill / Harness / Scaffold / 数据底座"])
    card(s, Inches(4.60), Inches(5.55), Inches(3.42), Inches(1.05), "contract",
         title="唯一主命题", lines=["P1 cost-aware separability"])
    card(s, Inches(8.25), Inches(5.55), Inches(3.42), Inches(1.05), "phys",
         title="判决规则", lines=["四态：支持 / 证伪 / 条件工程 / 不可判"])
    textbox(s, Inches(0.95), Inches(6.85), Inches(8.0), Inches(0.3),
            [("2026-08-22", 10, False, MUT)])
    return s


def s02_agenda(prs):
    s = blank(prs)
    heading(s, "本次评审要回答的三个问题", "议程与判据", tag="议程")
    items = [
        ("logic", "① 合回单篇是否正确？", [
            "v21/v22 拆分产生 4 处硬不一致：标题字面相同、契约元组字母冲突、",
            "随机化单位互相矛盾、两篇零交叉引用。合并可一次性消除全部 4 处。",
            "代价：放弃 v22 作为独立可投测量协议论文的形态。"]),
        ("contract", "② 四态判决是否成立？", [
            "六项义务是合取，任一不过阈则 P1 不可判 —— inconclusive 概率膨胀。",
            "企业篇原有 conditional-engineering 降级档，聚焦篇没有。",
            "本版升为全局第四态，使单个难测渠道不再永久锁死研究线。"]),
        ("phys", "③ 证据边界是否诚实？", [
            "本架构至今零实测结果，v22 与 v23 复核连续两轮标记为首要风险。",
            "评审需确认：是否接受以 research-program paper 形态推进，",
            "还是先跑单集群 2×2 可行性 pilot。"]),
    ]
    y = Inches(1.80)
    for pal, t, body in items:
        card(s, Inches(0.62), y, Inches(12.1), Inches(1.55), pal, title=t, lines=body)
        y += Inches(1.72)
    footer(s, 2, TOTAL, "判据均已落在 academy/v25-outline.md")
    return s


def s03_claim(prs):
    s = blank(prs)
    heading(s, "一句话主张与单一研究问题", "narrow contribution，不是泛化的三层架构", tag="主命题")
    card(s, Inches(0.62), Inches(1.82), Inches(12.1), Inches(1.30), "contract",
         title="研究问题（全文唯一）", title_size=15,
         lines=[("能否让 runtime 在激活可独立部署的能力时不实质改变容量响应，"
                 "同时在扩容时不实质改变能力语义，且执法成本可接受？", 13)])
    card(s, Inches(0.62), Inches(3.32), Inches(5.92), Inches(1.62), "logic",
         title="逻辑扩展 = 激活能力", lines=[
             "改变 activated capability configuration c",
             "= 在既有工作负载的准入路径上激活版本化 Skill bundle",
             "注册表里加不激活的条目 → 不是 treatment"])
    card(s, Inches(6.80), Inches(3.32), Inches(5.92), Inches(1.62), "phys",
         title="物理扩展 = 兼容扩容", lines=[
             "改变 Scaffold capacity configuration s",
             "= 改 worker 数 / 资源类别 / 隔离拓扑",
             "逻辑准入策略、模型设置、数据快照保持固定"])
    card(s, Inches(0.62), Inches(5.14), Inches(12.1), Inches(1.42), "neut",
         title="声明的边界", lines=[
             "① 这不是普适模块化，也不是证明某个实现可扩展；是对指定 runtime、workload、policy、"
             "数据快照与兼容容量区间的经验判断。",
             "② P1 不约束 Scaffold 主效应 —— 扩容让吞吐变好是预期的，不构成对 P1 的证伪。"])
    footer(s, 3, TOTAL, "来源：separability-study main.tex §1–§2")
    return s


def s04_objects(prs):
    s = blank(prs)
    heading(s, "四个责任对象", "责任边界，不是必须的团队或服务划分", tag="责任模型")
    specs = [
        ("logic", "Skill", "版本化行为", [
            "拥有任务意图、类型化输入输出、",
            "许可效果、前置条件、测试、发布身份",
            "不选物理放置、不自授凭证"]),
        ("contract", "Harness", "逻辑准入与控制", [
            "把请求 + 激活 bundle 编译为解析后契约",
            "路径构造、schema 校验、策略求值、",
            "效果授权、预算、版本身份、证据义务"]),
        ("phys", "Scaffold", "物理执行与隔离", [
            "算力、进程/容器边界、网络与数据局部性、",
            "运行时身份、资源计量、排队、故障隔离、attestation",
            "向 Harness 暴露资源事实，不重解释任务目标"]),
        ("data", "外置数据底座", "权威语义与快照", [
            "来源权威、语义定义、快照、血缘、",
            "凭证、留存、策略相关数据事实",
            "既不是隐藏的 Scaffold 库，也不是 Skill 私有缓存"]),
    ]
    x = Inches(0.62)
    w = Inches(2.94)
    for pal, name, role, body in specs:
        card(s, x, Inches(1.82), w, Inches(3.30), pal, title=name, title_size=17,
             lines=[(role, 11.5)] + body, body_size=9.8)
        x += w + Inches(0.16)
    card(s, Inches(0.62), Inches(5.34), Inches(12.1), Inches(1.10), "neut",
         lines=["一个实现可以把多个对象合并，但它的测量仍必须区分**逻辑准入、物理执行、数据权威**三者。",
                "信息隐藏只在被隐藏的决策拥有显式接口与归属时才有价值（Parnas 1972）。"])
    footer(s, 4, TOTAL, "契约元组本版统一为 C=⟨I,O,G,A,B,V⟩（原企业篇 ⟨I,O,G,B,E,T⟩ 已废弃）")
    return s


def s05_contract(prs):
    s = blank(prs)
    heading(s, "契约边界：C = ⟨I, O, G, A, B, V⟩", "解析后契约是可检查的最小单元；每个字段可独立被拒",
            tag="R2 统一")
    fields = [
        ("I", "类型化输入", "字段、取值、容差、口径"),
        ("O", "类型化输出", "结果子域 σ_out、schema"),
        ("G", "激活图", "节点与边、类型化衔接"),
        ("A", "授权与效果集", "authority、effects、证据义务"),
        ("B", "预算与绑定", "资源类别、上限、并发、隔离、驻留"),
        ("V", "版本钉定", "policy / model / data / verifier / trace"),
    ]
    x = Inches(0.62)
    w = Inches(1.95)
    for k, name, desc in fields:
        card(s, x, Inches(1.84), w, Inches(1.60), "contract", title=k, title_size=22,
             lines=[(name, 11, ), (desc, 9.2)], body_size=9.2)
        x += w + Inches(0.078)
    card(s, Inches(0.62), Inches(3.66), Inches(5.92), Inches(1.48), "neut",
         title="为什么这是关键", lines=[
             "契约让**选择性激活**可观测：注册表可以有很多 Skill，",
             "但只有准入路径上的激活行为改变时 c 才改变。",
             "这防止把注册表基数误当成能力干预，给实验一个稳定的处理身份。"])
    card(s, Inches(6.80), Inches(3.66), Inches(5.92), Inches(1.48), "warn",
         title="规划是概率的，准入不是", lines=[
             "Planner / executor 输出是**对授权的请求**，永不是授权本身。",
             "给定固定契约与策略快照，准入决策必须可重放；",
             "计划生成不必可重放。这是应对 prompt injection 的方式 ——",
             "靠契约解析，而不是靠指令过滤。"])
    card(s, Inches(0.62), Inches(5.36), Inches(12.1), Inches(1.06), "logic",
         title="有界组合（bounded composition）", lines=[
             "运行中的路径发现超出授权的需求时，当前 agent **不就地获得**该授权，而是返回类型化的未满足依赖；",
             "Harness 可解析一个新的有界子 agent 契约并作为子节点准入。终止取决于声明的后置条件，而非模型判断。"])
    footer(s, 5, TOTAL)
    return s


def s06_planes(prs):
    s = blank(prs)
    heading(s, "控制面 / 数据面切分", "画两个面不等于分离了 —— 完整中介与调度独立性义务必须观测每一次穿越",
            tag="架构")
    card(s, Inches(0.62), Inches(1.82), Inches(5.92), Inches(2.40), "contract",
         title="控制面 CP（∝ 状态事件）", lines=[
             "注册表、策略、准入、放置逻辑",
             "契约翻译、路由、维护子系统触发、不变量执行",
             "与业务流量解耦：QPS 上升不应同比抬高控制开销"])
    card(s, Inches(6.80), Inches(1.82), Inches(5.92), Inches(2.40), "phys",
         title="数据面 DP（∝ token 流量）", lines=[
             "只接收解析后契约，执行已准入的工作",
             "sandbox 计算、serving 前向、工具执行流",
             "控制规则复杂化不应抬高每请求 token 成本"])
    card(s, Inches(0.62), Inches(4.42), Inches(12.1), Inches(1.10), "warn",
         title="概率性控制面 —— 与经典 SDN/K8s 的真正差异", lines=[
             "经典控制面是确定性配置下发；本架构控制面部分由 LLM 驱动（契约翻译、反思、tool synthesis），",
             "控制决策本身不确定，因此需要可审计日志 + 不变量执行来约束其不确定性。"])
    card(s, Inches(0.62), Inches(5.64), Inches(12.1), Inches(0.86), "neut",
         lines=["可能的再耦合渠道：缓存填充、配额检查、重试控制器、放置提示 —— 任何绕过解析后契约的路径。"])
    footer(s, 6, TOTAL)
    return s


def s07_data(prs):
    s = blank(prs)
    heading(s, "数据子系统：功能契约 × 存储分层", "两个视角正交 —— 同一套 D₁–D₄ 可在任意记忆层被调用",
            tag="数据")
    card(s, Inches(0.62), Inches(1.82), Inches(5.92), Inches(2.72), "data",
         title="功能契约视角 𝒟 = ⟨D₁, D₂, D₃, D₄, Ω⟩", lines=[
             "D₁ on-policy 取数 API",
             "D₂ off-policy 语义总结平面（schema-on-read，反 ETL 截断）",
             "D₃ 数据治理 memory（沉淀 per-use-case data-usage skill）",
             "D₄ lifetime 体系（时新性 + 业务口径 + 访问 NFR 预算）"])
    card(s, Inches(6.80), Inches(1.82), Inches(5.92), Inches(2.72), "logic",
         title="存储分层视角 ℳ = ⟨Raw, Index, Theme⟩", lines=[
             "Tier 1 Raw — 物理着陆层，对接 Lakehouse",
             "Tier 2 Index — 语义编译层，OPIC 离线集群自底向上构建",
             "Tier 3 Theme — 经验沉淀层，自顶向下记录“数据被怎么用”",
             "Theme 是纯控制面：决定怎么取数，不搬运数据 → 支持平面分离"])
    card(s, Inches(0.62), Inches(4.72), Inches(12.1), Inches(1.02), "contract",
         title="IR 中间关系层 + 5W1H+Which 七维", lines=[
             "把契约映射里含糊的 “plan” 步骤实体化为 theme → 数据源集合的显式关系记录；",
             "除 semantic join 命中集外强制携带七维：时效性 / 使用范围 / 数据拥有者 / 数据语义 / 集成逻辑 / 物理访问方式 / 关系上下文（Which）。"])
    card(s, Inches(0.62), Inches(5.86), Inches(12.1), Inches(0.72), "neut",
         lines=["编号提示：三阶记忆命题在笔记侧已由 P10–P12 重编为 P18–P20（让位给并行度维度），且**不进正文命题表**。"])
    footer(s, 7, TOTAL, "来源：笔记 06 (v0.2) / 08-数据层修订 / 17-DataWiki-ThemeWiki-IR")
    return s


def s08_obligations(prs):
    s = blank(prs)
    heading(s, "六项测量义务", "不是靠检视授予的架构美德，而是每条都有五项必报输出", tag="可证伪")
    obs = [
        ("typed closure 类型化闭包", "执行或请求了 bundle 之外的操作"),
        ("complete mediation 完整中介", "调用到达效果边界而无当前契约授权"),
        ("effect non-interference 效果不干扰", "一次运行改变了另一次的授权效果或语义输入"),
        ("shared-state isolation 共享状态隔离", "越出契约分区的访问或未声明的争用边"),
        ("resource invariance 资源不变性", "改变兼容 s 时改变了未声明的语义输入"),
        ("scheduler independence 调度独立性", "未声明的调度→语义或能力→放置反馈边"),
    ]
    y = Inches(1.80)
    for i, (name, viol) in enumerate(obs):
        pal = "logic" if i < 2 else ("contract" if i < 4 else "phys")
        card(s, Inches(0.62), y, Inches(7.55), Inches(0.66), pal,
             lines=[(f"{name}   —   违规 = {viol}", 10.5)], body_size=10.5)
        y += Inches(0.735)
    card(s, Inches(8.42), Inches(1.80), Inches(4.30), Inches(2.10), "neut",
         title="五项必报输出", lines=[
             "① 仪器覆盖率  ② 观测到的违规",
             "③ 不确定性  ④ 执法成本",
             "⑤ 运行域 Ω 的排除项"])
    card(s, Inches(8.42), Inches(4.02), Inches(4.30), Inches(2.02), "warn",
         title="三个预设阈值", lines=[
             "覆盖率下限 γ、违规率上限 ν、监测灵敏度下限 η",
             "未覆盖事件一律**按违规计**（保守方向）",
             "盲化诊断注入校准灵敏度；",
             "自报的 Harness trace 不能自证完备"])
    footer(s, 8, TOTAL, "通过六项阈值不证明不存在第七条耦合渠道，只证明预登记威胁模型被测得足够好")
    return s


def s09_study(prs):
    s = blank(prs)
    heading(s, "Cluster-Period 随机交叉试验", "本版唯一的随机化单位 —— 企业篇原有的 “one independent run” 已删除",
            tag="R1 统一")
    card(s, Inches(0.62), Inches(1.82), Inches(3.90), Inches(2.16), "contract",
         title="实验单位", lines=[
             "cluster-period（system epoch）：",
             "隔离的 worker 池 / 租户分区 / 部署单元，",
             "在一个指派的 (c, s) 下运行固定周期。",
             "周期内的请求是**重复观测**，不是独立随机化单位。"])
    card(s, Inches(4.68), Inches(1.82), Inches(3.90), Inches(2.16), "logic",
         title="指派与平衡", lines=[
             "集群穿越所有合格的 (c, s) 组合；",
             "在集群与失效模式区块内随机化；",
             "用对抗平衡序列使每个处理出现在早晚位置相当，",
             "以便估计差异性结转。"])
    card(s, Inches(8.83), Inches(1.82), Inches(3.90), Inches(2.16), "phys",
         title="重置哨兵", lines=[
             "排空在途工作、恢复快照、清缓存与会话、",
             "重置限流与重试、重初始化调度历史。",
             "四项哨兵全部通过才能进入下一周期。"])
    card(s, Inches(0.62), Inches(4.18), Inches(5.92), Inches(1.32), "neut",
         title="三个观测量与判据", lines=[
             "R(c,s) 运行时响应 —— 主端点 Y_R = log 周期级 p95 延迟；m_R = log(1.10) 双侧 TOST 等效",
             "Q(c,s) 语义结局 —— 全需求满足占比；m_Q = 0.05 单侧非劣效"])
    card(s, Inches(6.80), Inches(4.18), Inches(5.92), Inches(1.32), "warn",
         title="E(c,s) 的边界（本版新增声明）", lines=[
             "E 估计的是**边际运行时执法开销**。配对回放固定 model/tool 输出，",
             "故“准入拒绝改变模型后续轨迹”的机会成本不计入 E，而归入 Q(c,s)。",
             "不声明这一点，会被读成低估了执法成本。"])
    card(s, Inches(0.62), Inches(5.66), Inches(12.1), Inches(0.80), "neut",
         lines=["待决：重置哨兵的 ±log(1.05) 与 −0.025 目前是裸数字。本版规则要求每个 margin 由具名决策推导并记录 owner —— 这两个数字必须补论证或标注为占位。"])
    footer(s, 9, TOTAL)
    return s


def s10_decision(prs):
    s = blank(prs)
    heading(s, "四态判决规则", "本版最大的结构变化：conditional-engineering 从企业篇专有升为全局第四态",
            tag="R3 新增", tagpal="phys")
    states = [
        ("data", "supported within Ω", [
            "六项义务全部通过覆盖率、灵敏度校准与最坏情况违规上界；",
            "90% 交互区间落在 [−m_R, m_R] 内；",
            "每个语义下界超过 −m_Q；每项执法开销上界低于预算。"]),
        ("warn", "falsified within Ω", [
            "合格性与校准足以判决，且交互区间整体越出 margin，",
            "或语义上界低于 −m_Q，或观测违规下界超过上限，",
            "或执法开销下界超出预算。"]),
        ("phys", "conditional-engineering ★", [
            "一项或多项义务违规率超阈 → 报告完整结局估计 + 越阈条件及其观测值，",
            "**对猜想不作任何方向的判决**。",
            "它刻画的是 runtime 的工程状态，而非 P1 的科学地位；可发表且有信息量。"]),
        ("neut", "inconclusive", [
            "仪器覆盖或监测校准不足、区间无法分辨 margin、",
            "重置后仍有污染、预声明合格性规则失败，",
            "或所需量无法以计划的整群不确定性估计。"]),
    ]
    y = Inches(1.80)
    for pal, name, body in states:
        card(s, Inches(0.62), y, Inches(12.1), Inches(1.16), pal, title=name,
             title_size=14, lines=body, body_size=10)
        y += Inches(1.28)
    footer(s, 10, TOTAL, "为什么必要：六项义务是合取，任一不过阈就锁死判决 —— 第四态让单个难测渠道（如调度独立性）不再永久锁死研究线")
    return s


def s11_fixed(prs):
    s = blank(prs)
    heading(s, "合并直接修掉的跨篇不一致", "上一轮审查的 P0/P1 级问题，逐条对应", tag="合并收益")
    rows = [
        ("P0-1", "两篇标题字面完全相同", "latex_to_preprint.py:29 单一硬编码 TITLE，无 --title 参数",
         "合并后自动消失；仍补 --title 防复发"),
        ("P0-2", "契约元组字母冲突", "聚焦篇 ⟨I,O,G,A,B,V⟩ vs 企业篇 ⟨I,O,G,B,E,T⟩",
         "全局统一为 ⟨I,O,G,A,B,V⟩，字段级 walkthrough 逐字段改写"),
        ("P0-3", "随机化单位互相矛盾", "企业篇 “one independent run” vs 聚焦篇 “非独立随机化单位”",
         "以 cluster-period 为准，删除企业篇整句"),
        ("P0-4", "两篇零交叉引用", "无 Part I / Part II / companion 任何互锚",
         "合并后自动消失"),
        ("P1-5", "降级档只做了一半", "conditional-engineering 仅企业篇有，聚焦篇三态无此档",
         "升为全局第四态"),
        ("P1-6", "margin 推导规则不对称", "企业篇立了具名决策规则，聚焦篇哨兵是裸数字",
         "规则全局化，裸数字补论证或标占位"),
        ("P1-7", "E(c,s) 未声明边际性", "全文无 marginal / opportunity cost 相关声明",
         "在 E 定义处显式声明，机会成本归 Q"),
    ]
    y = Inches(1.76)
    for tag, what, evid, fix in rows:
        pal = "warn" if tag.startswith("P0") else "phys"
        card(s, Inches(0.62), y, Inches(1.00), Inches(0.70), pal,
             lines=[(tag, 11)], body_size=11)
        textbox(s, Inches(1.76), y + Inches(0.02), Inches(3.55), Inches(0.68),
                [(what, 10.5, True, INK)])
        textbox(s, Inches(5.40), y + Inches(0.03), Inches(3.60), Inches(0.66),
                [(evid, 9, False, MUT)])
        textbox(s, Inches(9.15), y + Inches(0.03), Inches(3.58), Inches(0.66),
                [("→ " + fix, 9.2, False, C("data", 2))])
        y += Inches(0.76)
    footer(s, 11, TOTAL, "另有 P2-8/9 文献覆盖、P3-10 篇幅去重、P4-11/12/13 笔记单一事实源，见 v25-outline.md")
    return s


def s12_biblio(prs):
    s = blank(prs)
    heading(s, "文献覆盖的空洞", "26 篇已下载 PDF 中 9 篇零引用 —— 其中一篇构成 novelty boundary 的实质风险",
            tag="P2-8", tagpal="warn")
    card(s, Inches(0.62), Inches(1.82), Inches(12.1), Inches(1.42), "warn",
         title="最危险的一篇：arXiv 2605.26112", title_size=15, lines=[
             "《From Model Scaling to System Scaling: From Training to Serving of LLM-Based Agentic Systems》",
             "题目即“模型扩展 → 系统扩展”，与本文 capability / capacity 双轴属同一问题域。",
             "vault 里已有专门笔记，但两篇论文的 bib 都没有它 —— 审稿人极可能据此质疑 novelty。"])
    items = [
        ("contract", "2607.13987", "Agent Skill Security：威胁模型、攻击、防御与评测", "§3 威胁模型的直接对标，semi-trusted Skills 一条尤其需要"),
        ("phys", "2605.22781", "DeltaBox：LLM agent 沙箱的轻量容器", "Scaffold 隔离基底的具体实现，支撑资源不变性义务"),
        ("logic", "2605.18747", "Code as Agent Harness 综述", "Harness 综述，定位 Skill-as-Code；vault 已有笔记"),
    ]
    y = Inches(3.44)
    for pal, aid, title, why in items:
        card(s, Inches(0.62), y, Inches(12.1), Inches(0.82), pal,
             lines=[(f"{aid}   {title}", 11), (why, 9.4)], body_size=9.4)
        y += Inches(0.94)
    card(s, Inches(0.62), Inches(6.28), Inches(12.1), Inches(0.62), "neut",
         lines=["引用卫生：标题与完整作者列表必须经 arXiv API 核实后才写入 references.bib（本仓库有伪造引用元数据的历史，CHANGELOG v4 修 8 条、v8 修 3 条）。"])
    footer(s, 12, TOTAL)
    return s


def s13_threat(prs):
    s = blank(prs)
    heading(s, "威胁模型", "每条假设都是架构可被攻击的地方，而不是它建立的性质", tag="安全")
    rows = [
        ("contract", "可信计算基", "控制面可信：Harness 编译器、Skill 注册表、策略求值器、gate 实现、Scaffold attestation。任一被攻破则边界失效。"),
        ("warn", "不可信模型输出", "Planner/executor 输出是对授权的请求，永不是授权。计划命名契约外操作 = 被拒请求，不是被放宽的契约。"),
        ("phys", "半可信 Skill", "作者已认证且可追责，但不假设正确或善意。Skill 可以声明一种行为而表现另一种，故声明是待行为性验证的准入输入。"),
        ("logic", "不可信外部内容", "工具结果、检索文档、数据源载荷均可被攻击者影响。数据面文本不能改变激活路径、授权效果集或策略版本。"),
        ("neut", "对手能力", "可以撰写/修改 Skill、影响请求文本、控制外部内容、把已批准 Skill 组合成危险路径；不能破坏 Scaffold 隔离、伪造 attestation 或颠覆控制面。"),
    ]
    y = Inches(1.80)
    for pal, name, body in rows:
        card(s, Inches(0.62), y, Inches(12.1), Inches(0.90), pal, title=name,
             title_size=12.5, lines=[(body, 9.8)], body_size=9.8)
        y += Inches(1.00)
    card(s, Inches(0.62), Inches(6.86), Inches(12.1), Inches(0.42), "neut",
         lines=[("目标很窄：对手能达到的每个效果都必须出现在某个解析后契约的授权效果集中，"
                 "集合外的每次尝试都是被记录的类型化失败。gate 与 trace 使违规可归因，但都不阻止已授权效果造成危害。", 9.4)])
    footer(s, 13, TOTAL)
    return s


def s14_evidence(prs):
    s = blank(prs)
    heading(s, "证据边界：本架构至今没有实测结果", "v22 与 v23 复核连续两轮标记为首要风险 —— 评审必须知道这一点",
            tag="必读", tagpal="warn")
    card(s, Inches(0.62), Inches(1.82), Inches(12.1), Inches(1.24), "warn",
         title="本文不提供什么", title_size=15, lines=[
             "没有完成的 runtime 实现、没有实测的执法预算、没有数据集、没有 benchmark 结果、",
             "没有关于某个部署可支撑规模的主张。图是责任与协议视图，不是截图或实现证据。"])
    rows = [
        ("supported 不等于“证明独立”", "它只表示在预登记的 Ω、margin、功效与仪器条件内未触发证伪标准。任何更强外推都不成立。"),
        ("六条件合取可能贵到不可测", "论文自承：严格的覆盖规则降低假支持，但会让有用的部署变成 inconclusive。外部效度因此存疑。"),
        ("组织与治理实证薄弱", "四类责任对象的 owner 划分、争议升级（现为规则草案）与跨企业可迁移性，尚未经组织研究验证。"),
        ("唯一剩余的 P1 级评审项", "与 LangGraph 等现有系统的逐项对比表，自 v20 评审第 6 条提出至 v23 仍未落地。"),
    ]
    y = Inches(3.26)
    for name, body in rows:
        card(s, Inches(0.62), y, Inches(12.1), Inches(0.80), "neut", title=name,
             title_size=12, lines=[(body, 9.6)], body_size=9.6)
        y += Inches(0.90)
    footer(s, 14, TOTAL, "这一页不得从评审材料中删去 —— 隐去它会误导评审")
    return s


def s15_next(prs):
    s = blank(prs)
    heading(s, "下一步：最便宜的第一步是 pilot，不是继续扩写", "目的不是检验 P1，而是验证协议可执行",
            tag="行动")
    card(s, Inches(0.62), Inches(1.82), Inches(5.92), Inches(2.50), "data",
         title="单集群 2×2 序列可行性 pilot", title_size=15, lines=[
             "固定模型、policy 与 workload family，在两个 Scaffold class 上",
             "跑一轮 cluster-period 随机交叉。要验证的是三件事：",
             "① 重置哨兵在真实 runtime 上可达成",
             "② 六项义务的覆盖率下限 γ 可达",
             "③ 配对安全回放可实施",
             "无论结果是 supported / falsified / conditional / inconclusive，",
             "都比继续扩写架构说明更能提高可信度。"])
    card(s, Inches(6.80), Inches(1.82), Inches(5.92), Inches(2.50), "contract",
         title="论文侧本轮之后的工作", title_size=15, lines=[
             "① 逐段合并 822 + 308 行 LaTeX，重排为单套 8 图",
             "② 按 R1–R4 落实四条硬规则",
             "③ 补 4 篇必引文献并重写 novelty boundary",
             "④ 重写 notes/07 创新点 ledger（当前仍是 v0.1 / 2026-06-20）",
             "⑤ 补 LangGraph 对比表",
             "⑥ 给 latex_to_preprint.py 加 --title 参数"])
    card(s, Inches(0.62), Inches(4.52), Inches(12.1), Inches(1.00), "warn",
         title="需要评审拍板的两件事", lines=[
             "① 是否接受“放弃 v22 作为独立可投测量协议论文”这一代价？（合回单篇的直接后果）",
             "② 是否先跑 pilot 再投稿，还是以 research-program paper 形态先投？"])
    card(s, Inches(0.62), Inches(5.72), Inches(12.1), Inches(0.70), "neut",
         lines=["笔记侧已完成：三份漂移副本合并为单一权威目录，命题编号撞车（P10–P12 两套）已修，repo notes/ 降为镜像。"])
    footer(s, 15, TOTAL)
    return s


def s16_end(prs):
    s = blank(prs)
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.22), SH)
    band.fill.solid()
    band.fill.fore_color.rgb = C("contract", 1)
    band.line.fill.background()
    band.shadow.inherit = False
    textbox(s, Inches(0.95), Inches(2.70), Inches(11.2), Inches(1.0),
            [("讨论", 40, True, INK)])
    textbox(s, Inches(0.95), Inches(3.78), Inches(11.2), Inches(1.6),
            [("计划    docs/superpowers/plans/2026-08-22-v25-remerge-and-review-deck.md", 12, False, MUT),
             ("大纲    academy/v25-outline.md", 12, False, MUT),
             ("本 deck  academy/architecture-review-v25-zh.pptx（构建脚本 review_deck_build.py）", 12, False, MUT),
             ("笔记    Obsidian / Agentic-Runtime-参考架构（唯一权威，repo notes/ 为镜像）", 12, False, MUT)])
    return s


def main():
    prs = Presentation()
    prs.slide_width, prs.slide_height = SW, SH
    for fn in (s01_cover, s02_agenda, s03_claim, s04_objects, s05_contract,
               s06_planes, s07_data, s08_obligations, s09_study, s10_decision,
               s11_fixed, s12_biblio, s13_threat, s14_evidence, s15_next, s16_end):
        fn(prs)
    out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                       "architecture-review-v25-zh.pptx")
    prs.save(out)
    print(f"wrote {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
