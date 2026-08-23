#!/usr/bin/env python3
"""论文报告 deck 构建器 — 对应合并稿 v27。

产物: academy/paper-report-v27-zh.pptx (16:9, 中文)
用法: python3 academy/paper_report_build.py

与 review_deck_build.py 的分工：
  - review deck  : 架构评审用，重心在责任边界、契约、判决规则的取舍
  - 本 deck      : 论文报告用，按论文章节顺序讲，并内嵌论文的 7 张真实插图
    （agentic-runtime-preprint/paper_source/{1..7}.png，与 PDF 中所见一致），
    论文没有配图的部分（P1 端点、六义务、四态判决、v27 形式化附录）改为矢量绘制。

配色沿用 site/style.css 迁移后的 Tailwind 系。
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

SW, SH = Inches(13.333), Inches(7.5)
YAHEI = "Microsoft YaHei"
HERE = os.path.dirname(os.path.abspath(__file__))
FIGDIR = os.path.join(HERE, "agentic-runtime-preprint", "paper_source")
FIG_ASPECT = 1707 / 952

PAL = {
    "logic":    ("EFF6FF", "3B82F6", "1E40AF"),
    "contract": ("F5F3FF", "7C3AED", "5B21B6"),
    "phys":     ("FFF7ED", "EA580C", "C2410C"),
    "data":     ("ECFDF5", "059669", "047857"),
    "neut":     ("F9FAFB", "E5E7EB", "374151"),
    "warn":     ("FEF2F2", "DC2626", "991B1B"),
}
INK = RGBColor(0x11, 0x18, 0x27)
MUT = RGBColor(0x6B, 0x72, 0x80)
PAGE = RGBColor(0xFF, 0xFF, 0xFF)
TOTAL = 19


def C(name, i):
    return RGBColor.from_string(PAL[name][i])


def _yahei(run):
    run.font.name = YAHEI
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {})
            rPr.append(e)
        e.set("typeface", YAHEI)


def blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    f = s.background.fill
    f.solid()
    f.fore_color.rgb = PAGE
    return s


def textbox(slide, x, y, w, h, lines, *, size=11, color=INK, bold=False,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        p.space_after = Pt(3)
        txt, sz, bd, col = ln if isinstance(ln, tuple) else (ln, size, bold, color)
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(sz)
        r.font.bold = bd
        r.font.color.rgb = col
        _yahei(r)
    return tb


def heading(slide, title, sub=None, *, tag=None, tagpal="contract"):
    textbox(slide, Inches(0.62), Inches(0.36), Inches(10.5), Inches(0.60),
            [(title, 25, True, INK)])
    if sub:
        textbox(slide, Inches(0.62), Inches(0.96), Inches(11.5), Inches(0.40),
                [(sub, 11.5, False, MUT)])
    if tag:
        w = Inches(0.30 + 0.115 * len(tag))
        s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   SW - Inches(0.62) - w, Inches(0.40), w, Inches(0.33))
        s.fill.solid()
        s.fill.fore_color.rgb = C(tagpal, 0)
        s.line.color.rgb = C(tagpal, 1)
        s.line.width = Pt(1.0)
        s.shadow.inherit = False
        try:
            s.adjustments[0] = 0.22
        except Exception:
            pass
        tf = s.text_frame
        tf.word_wrap = False
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = tag
        r.font.size = Pt(10)
        r.font.bold = True
        r.font.color.rgb = C(tagpal, 2)
        _yahei(r)
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.62), Inches(1.42),
                                SW - Inches(1.24), Pt(1.1))
    ln.fill.solid()
    ln.fill.fore_color.rgb = RGBColor.from_string("E5E7EB")
    ln.line.fill.background()
    ln.shadow.inherit = False


def card(slide, x, y, w, h, pal, *, title=None, lines=None, title_size=13,
         body_size=10, dash=False):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = C(pal, 0)
    s.line.color.rgb = C(pal, 1)
    s.line.width = Pt(1.2)
    s.shadow.inherit = False
    if dash:
        el = s.line._get_or_add_ln()
        el.append(el.makeelement(qn("a:prstDash"), {"val": "dash"}))
    try:
        s.adjustments[0] = 0.06
    except Exception:
        pass
    tf = s.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.09)
    tf.margin_bottom = Inches(0.07)
    first = True
    if title:
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = title
        r.font.size = Pt(title_size)
        r.font.bold = True
        r.font.color.rgb = C(pal, 2)
        _yahei(r)
        first = False
    for ln in (lines or []):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.13
        p.space_before = Pt(3)
        txt, sz = (ln, body_size) if isinstance(ln, str) else ln
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(sz)
        r.font.color.rgb = INK
        _yahei(r)
    return s


def figure(slide, n, x, y, w):
    """内嵌论文第 n 张插图（与 PDF 中所见同一文件）。"""
    path = os.path.join(FIGDIR, f"{n}.png")
    if not os.path.exists(path):
        raise FileNotFoundError(path)
    h = Inches(w.inches / FIG_ASPECT)
    frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   x - Inches(0.06), y - Inches(0.06),
                                   w + Inches(0.12), h + Inches(0.12))
    frame.fill.solid()
    frame.fill.fore_color.rgb = RGBColor.from_string("F9FAFB")
    frame.line.color.rgb = RGBColor.from_string("E5E7EB")
    frame.line.width = Pt(1.0)
    frame.shadow.inherit = False
    try:
        frame.adjustments[0] = 0.03
    except Exception:
        pass
    slide.shapes.add_picture(path, x, y, width=w)
    textbox(slide, x, y + h + Inches(0.10), w, Inches(0.26),
            [(f"论文 Figure {n}（原图，非重绘）", 8.5, False, MUT)])
    return h


def footer(slide, n, note=""):
    textbox(slide, Inches(0.62), SH - Inches(0.46), Inches(9.4), Inches(0.28),
            [(note, 8.5, False, MUT)])
    textbox(slide, SW - Inches(1.55), SH - Inches(0.46), Inches(0.95), Inches(0.28),
            [(f"{n} / {TOTAL}", 8.5, False, MUT)], align=PP_ALIGN.RIGHT)


def figure_slide(prs, idx, fig_n, title, sub, tag, tagpal, cards, note=""):
    """左图右文的标准版式。"""
    s = blank(prs)
    heading(s, title, sub, tag=tag, tagpal=tagpal)
    fw = Inches(8.05)
    figure(s, fig_n, Inches(0.68), Inches(1.72), fw)
    x = Inches(9.05)
    w = Inches(3.66)
    top, bottom = 1.72, 6.94          # 6.94 leaves the footer clear
    gap = 0.14
    heights = [h for _, _, _, h in cards]
    budget = (bottom - top) - gap * (len(cards) - 1)
    if sum(heights) > budget:         # 等比压缩而不是逐页手调
        k = budget / sum(heights)
        heights = [h * k for h in heights]
    y = top
    for (pal, t, body, _), h in zip(cards, heights):
        card(s, x, Inches(y), w, Inches(h), pal, title=t, lines=body)
        y += h + gap
    footer(s, idx, note)
    return s


# ─────────────────────────── slides ───────────────────────────

def s01(prs):
    s = blank(prs)
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.22), SH)
    band.fill.solid()
    band.fill.fore_color.rgb = C("contract", 1)
    band.line.fill.background()
    band.shadow.inherit = False
    textbox(s, Inches(0.95), Inches(1.95), Inches(11.2), Inches(0.44),
            [("论文报告 · 合并稿 v27", 14, True, C("contract", 2))])
    textbox(s, Inches(0.95), Inches(2.44), Inches(11.5), Inches(1.7),
            [("A Contract-Centered Architecture for", 33, True, INK),
             ("Scalable and Manageable Agentic Runtimes", 33, True, INK)])
    textbox(s, Inches(0.95), Inches(4.22), Inches(11.2), Inches(0.9),
            [("契约中心的可扩展 · 可管理 Agentic Runtime 参考架构", 13, False, MUT),
             ("Yaxiao Liu · 37 页 · 8 图 · 45 条引用 · 无实测结果（research-program paper）", 12, False, MUT)])
    card(s, Inches(0.95), Inches(5.42), Inches(3.42), Inches(1.16), "contract",
         title="唯一主命题", lines=["P1 cost-aware", "capability-capacity separability"])
    card(s, Inches(4.60), Inches(5.42), Inches(3.42), Inches(1.16), "logic",
         title="六项测量义务", lines=["γ 覆盖率 / ν 违规率 / η 灵敏度", "未覆盖事件按违规计"])
    card(s, Inches(8.25), Inches(5.42), Inches(3.42), Inches(1.16), "phys",
         title="四态判决", lines=["supported / falsified /", "conditional-engineering / inconclusive"])
    textbox(s, Inches(0.95), Inches(6.80), Inches(8.0), Inches(0.3),
            [("2026-08-23", 10, False, MUT)])
    return s


def s02(prs):
    s = blank(prs)
    heading(s, "论文要回答的一个问题", "全文只有一个研究问题，其余都是为使它可判而存在", tag="研究问题")
    card(s, Inches(0.62), Inches(1.74), Inches(12.1), Inches(1.24), "contract",
         title="Research question", title_size=14,
         lines=[("Can a runtime activate independently deployable capabilities without materially "
                 "changing capacity response, while scaling capacity without materially changing "
                 "capability semantics, at an acceptable enforcement cost?", 12.5)])
    card(s, Inches(0.62), Inches(3.14), Inches(5.92), Inches(1.62), "logic",
         title="产品方的变更请求", lines=[
             "想激活一个新行为，而不必重新设计执行机队。",
             "改变 activated capability configuration c ——",
             "注册表里加不激活的条目不算 treatment。"])
    card(s, Inches(6.80), Inches(3.14), Inches(5.92), Inches(1.62), "phys",
         title="平台方的变更请求", lines=[
             "想加 worker、隔离区、加速器或区域容量，",
             "而不改变已准入行为的含义。",
             "改变 Scaffold capacity configuration s。"])
    card(s, Inches(0.62), Inches(4.94), Inches(12.1), Inches(1.02), "neut",
         title="为什么单轴评测答不了", lines=[
             "能力工作用任务成功率评，基础设施工作用吞吐与延迟评 —— 两种单轴评测都看不出两轴是否经由共享状态、",
             "资源敏感输入、调度反馈、未声明副作用或绕过执法的路径**重新耦合**。"])
    card(s, Inches(0.62), Inches(6.10), Inches(12.1), Inches(0.72), "warn",
         lines=[("注意：P1 不约束 Scaffold 主效应。扩容让吞吐变好是预期的，不构成对 P1 的支持，也不构成证伪。", 10.5)])
    footer(s, 2, "论文 §1–§2")
    return s


def s03(prs):
    return figure_slide(
        prs, 3, 1, "四个责任对象", "责任边界，不是必须的团队或服务划分", "责任模型", "logic",
        [("logic", "Skill", ["版本化、可独立部署的行为声明。",
                             "拥有任务意图、类型化 I/O、许可效果、",
                             "前置条件、测试、发布身份。"], 1.42),
         ("contract", "Harness", ["运行时编译器与治理者。",
                                  "拥有逻辑准入与控制：路径构造、",
                                  "schema 校验、策略求值、效果授权。"], 1.42),
         ("phys", "Scaffold", ["物理执行与隔离边界。",
                               "三模块：Sandbox 环境 / API 双轨路由 / 治理。",
                               "向 Harness 暴露资源事实，不重解释任务目标。"], 1.52),
         ("data", "栈外数据底座", ["来源权威、语义定义、快照、血缘。",
                                   "既非隐藏的 Scaffold 库，",
                                   "也非 Skill 私有缓存。"], 1.42)],
        "论文 §5.1 · SkillMD-138K 证据：98.73% 有 frontmatter，但仅约 23.2% 声明可执行组件 → 闭包必须由 Harness 在准入时强制")


def s04(prs):
    return figure_slide(
        prs, 4, 2, "契约边界 C = ⟨I, O, G, A, B, V⟩", "解析后契约是可检查的最小单元；每个字段可独立被拒",
        "契约", "contract",
        [("contract", "六个字段", ["I / O 类型化输入输出",
                                   "G 激活图",
                                   "A 授权与效果集 + 证据义务",
                                   "B 预算与绑定约束",
                                   "V 版本钉定（policy/model/data/verifier/trace）"], 1.80),
         ("logic", "选择性激活可观测", ["注册表可以有很多 Skill，",
                                        "但只有准入路径上的激活行为改变时 c 才改变 ——",
                                        "这防止把注册表基数误当成能力干预。"], 1.52),
         ("warn", "规划是概率的，准入不是", ["Planner 输出是**对授权的请求**，永不是授权。",
                                             "应对 prompt injection 靠契约解析，",
                                             "不靠指令过滤。"], 1.42)],
        "论文 §6.1 · v25 统一了元组字母（企业篇原用 ⟨I,O,G,B,E,T⟩，已废弃）")


def s05(prs):
    return figure_slide(
        prs, 5, 3, "控制面 / 数据面", "画两个面不等于分离了 —— 义务必须观测每一次穿越", "架构", "logic",
        [("logic", "控制面 CP", ["注册表、策略、准入、放置逻辑。",
                                 "∝ 状态事件，与业务流量解耦。"], 1.16),
         ("phys", "数据面 DP", ["只接收解析后契约，执行已准入工作。",
                                "∝ token 流量。"], 1.16),
         ("warn", "控制态泄漏", ["把每个 Skill schema 与策略注入每次请求，",
                                 "会把逻辑扩展变成线性 prompt 增长，",
                                 "并产生**含糊的权限**。"], 1.42),
         ("neut", "概率性控制面", ["经典 SDN/K8s 控制面是确定性配置下发；",
                                   "此处部分由 LLM 驱动，故需可审计日志",
                                   "与不变量执行来约束不确定性。"], 1.42)],
        "论文 §6.2 · 可能的再耦合渠道：缓存填充、配额检查、重试控制器、放置提示")


def s06(prs):
    return figure_slide(
        prs, 6, 4, "派生闭包编排", "请求内的能力增长靠「加有界子 agent」，而非扩大既有 agent", "协议", "contract",
        [("contract", "闭包规则", ["每个子 agent 契约在命名操作、证据来源、",
                                   "I/O 与交接 schema、效果、预算上闭合。",
                                   "闭包外的引用**不就地满足**。"], 1.52),
         ("logic", "引用即派生条件", ["超出授权的需求记为**类型化未满足依赖**，",
                                      "触发 Harness 解析新契约并作为子节点准入。",
                                      "每次派生都过同一套确定性 gate。"], 1.52),
         ("data", "主 agent = 满足度控制器", ["不执行任务工具，只做 summarize + semantic join",
                                              "+ top-k，并按声明的 σ_out 子域算 satisfaction ratio。",
                                              "verifier 若不可校准，比值只是估计，不能单独授权终止。"], 1.62)],
        "论文 §7.4 · 终止判据是契约级谓词，而非模型判断")


def s07(prs):
    return figure_slide(
        prs, 7, 5, "数据底座：以源证据为权威", "v26 的实质修订 —— 摘要与索引降为可审计衍生物", "v26 修订", "data",
        [("data", "权威事实", ["canonical raw / source-native records、",
                               "版本快照、来源、访问策略、",
                               "确定性读取与抽取操作。"], 1.42),
         ("warn", "衍生物（不得替代原始证据）", ["摘要、索引、候选关系。",
                                                 "此前把它们当作可替代表示，v26 已改。"], 1.16),
         ("logic", "访问路径", ["全局排序发现 → 结构化导航 / 操作符",
                                "→ 读取版本化源证据。",
                                "在线时延与索引、导航的摊销成本分开记账。"], 1.42),
         ("neut", "被移除的实现性假设", ["「永久保留」「摘要即隐私保护」「SQLite 单写者」",
                                         "现由治理契约决定，不再写成架构性质。"], 1.30)],
        "论文 §8.1 · Memory 只记可观察工具轨迹与证据指针，不保存也不要求私有 chain-of-thought")


def s08(prs):
    s = blank(prs)
    heading(s, "中间关系层 IR 与 5W1H + Which", "两个正交 registry 之间唯一的显式耦合点", tag="数据", tagpal="data")
    card(s, Inches(0.62), Inches(1.72), Inches(3.86), Inches(1.36), "logic",
         title="Data Wiki", lines=["address / index / operator catalog：",
                                   "数据在哪、怎么寻址、有哪些确定性操作符。"])
    card(s, Inches(4.72), Inches(1.72), Inches(3.86), Inches(1.36), "phys",
         title="IR（唯一耦合点）", lines=["theme ↦ ⟨命中的源摘要集合, 5W1H+Which⟩",
                                          "记录来源、快照、操作符、证据、回退与兼容范围。"])
    card(s, Inches(8.82), Inches(1.72), Inches(3.90), Inches(1.36), "data",
         title="Theme Wiki", lines=["版本化任务族输出与质量契约：",
                                    "经训练验证后才登记的可复用产出模板。"])
    dims = [("When", "有效期与更新节奏", "logic"), ("Where", "访问域边界", "logic"),
            ("Who", "数据归属与授权", "contract"), ("What", "字段业务语义", "contract"),
            ("Why", "声明的集成理由 / 因果假设", "phys"), ("How", "物理访问路径", "phys"),
            ("Which", "关系上下文（新增第七维）", "data")]
    x = Inches(0.62)
    w = Inches(1.68)
    for name, desc, pal in dims:
        card(s, x, Inches(3.34), w, Inches(1.30), pal, title=name, title_size=13,
             lines=[(desc, 9)], body_size=9)
        x += w + Inches(0.05)
    card(s, Inches(0.62), Inches(4.86), Inches(12.1), Inches(1.06), "neut",
         title="为什么 Which 是必要的", lines=[
             "前六维只在孤立状态下限定一个 fact；Which 把它放进 semantic join 应当展开的可导航邻域 —— "
             "时间窗内共现的 temporal links 与 typed cross-references。",
             "缺 When 可能过期、缺 Who 可能越权、缺 How 不可达、缺 Which 则上下文孤立。同一份七维记录兼作 agent 侧记忆的索引原子。"])
    card(s, Inches(0.62), Inches(6.06), Inches(12.1), Inches(0.76), "warn",
         lines=[("P17 已在 v26 收窄为「有界策略契约解耦」：仅在兼容范围、界面复杂度、可替换性、"
                 "无隐藏共享状态与迁移规则均明确时成立。", 10.5)])
    footer(s, 8, "论文 §8.2 · 此节不提出新的检索或数据集成算法，贡献是把底座已有信息重组为可审计的关系层")
    return s


def s09(prs):
    return figure_slide(
        prs, 9, 7, "Skill 生命周期与冻结为 code", "draft → train → validate → staged release → monitor",
        "生命周期", "phys",
        [("phys", "操作闭包是验证 gate", ["测试须证明外部可见行为可由声明的效果表达。",
                                          "模型版本稳定性在同一边界测试。"], 1.30),
         ("contract", "确定性锚", ["当工具集稳定、输出结构化、分支可枚举，",
                                   "执行过程可冻结为 code，不再每轮概率生成 ——",
                                   "但只在输入、依赖快照、种子、执行语义固定时确定。"], 1.62),
         ("logic", "双子目标 Reward（P15）", ["r_out 沿 σ_out 子域独立标注 × r_proc 过程一致性。",
                                              "Pareto 式 gate：保护子域不退化 + 至少一个目标子域实质改善。"], 1.42),
         ("warn", "phantom violation 警告", ["有工作观测到 harness 优化器为**从未发生的失败**",
                                              "捏造 guardrail（60 次中 15 次）→ 过程判据必须是",
                                              "运行前登记的产物，不能在评分时由模型生成。"], 1.52)],
        "论文 §9 · P15 是次级假设，不进 P1 判决")


def s10(prs):
    s = blank(prs)
    heading(s, "P1：cost-aware capability-capacity separability", "唯一主命题 —— 三个端点、三个 margin",
            tag="主命题", tagpal="contract")
    card(s, Inches(0.62), Inches(1.70), Inches(3.90), Inches(2.10), "phys",
         title="运行时响应 R(c,s)", lines=[
             "主端点 Y_R = ln(周期级 p95 准入→终止延迟)",
             "估计量：差分中的差分 Δ_R",
             "margin：m_R = log(1.10)",
             "规则：90% 双侧整群置信区间整体落在",
             "[−m_R, m_R] 内（TOST，α=0.05）",
             "即 p95 尾延迟的交互不超过 10% 乘性变化"])
    card(s, Inches(4.72), Inches(1.70), Inches(3.90), Inches(2.10), "data",
         title="语义结局 Q(c,s)", lines=[
             "主端点 Q_req = 满足**每一条**预登记需求的运行占比",
             "对比：D_Q(c) = Q_req(c,s₁) − Q_req(c,s₀)",
             "margin：m_Q = 0.05",
             "规则：单侧 95% 整群置信下界 > −0.05",
             "（方向性非劣，非等效）"])
    card(s, Inches(8.82), Inches(1.70), Inches(3.90), Inches(2.10), "contract",
         title="执法开销 E(c,s)", lines=[
             "预算 B_E 在看到结局前固定，逐分量适用：",
             "时延 / 算力 / 证据量 / 金额上限",
             "每个分量的单侧 95% 上界须低于其预算",
             "廉价的未仪器化路径**不能**满足该预算 ——",
             "缺失的控制工作使证据不合格"])
    card(s, Inches(0.62), Inches(4.02), Inches(12.1), Inches(1.06), "warn",
         title="所有 margin 必须由具名决策推导", lines=[
             "不得取自 pilot 的观测方差 —— 那会让噪声大的仪器通过放宽自己的 margin 来给自己发证。",
             "决策、决策的 owner、结果数值须作为三个独立的预登记字段记录。reset 哨兵的两个数值本版**显式留空**，注册前必须补齐。"])
    card(s, Inches(0.62), Inches(5.22), Inches(12.1), Inches(0.80), "neut",
         lines=[("功效按交互项设计，而非按更大更容易的 Scaffold 主效应。工作样例：σ=0.4、ρ=0.05、n=500 时，"
                 "80% 功效下最小可检测交互约 0.08（8% 乘性），需每处理顺序约 41 个集群。", 10.5)])
    card(s, Inches(0.62), Inches(6.16), Inches(12.1), Inches(0.66), "warn",
         lines=[("把功效不足报成独立性结论，是这条研究线最可能产生「有利于自己猜想的假阳性」的方式。", 10.5)])
    footer(s, 10, "论文 §5.3")
    return s


def s11(prs):
    s = blank(prs)
    heading(s, "六项测量义务", "不是靠检视授予的架构美德 —— 每条五项必报输出 + 三阈值", tag="可证伪", tagpal="logic")
    obs = [("typed closure", "类型化闭包", "执行或请求了 bundle 之外的操作", "logic"),
           ("complete mediation", "完整中介", "调用到达效果边界而无当前契约授权", "logic"),
           ("effect non-interference", "效果不干扰", "一次运行改变了另一次的授权效果或语义输入", "contract"),
           ("shared-state isolation", "共享状态隔离", "越出契约分区的访问或未声明的争用边", "contract"),
           ("resource invariance", "资源不变性", "改变兼容 s 时改变了未声明的语义输入", "phys"),
           ("scheduler independence", "调度独立性", "未声明的调度→语义或能力→放置反馈边", "phys")]
    y = Inches(1.70)
    for en, zh, viol, pal in obs:
        card(s, Inches(0.62), y, Inches(7.60), Inches(0.63), pal,
             lines=[(f"{en}  ·  {zh}   —   违规 = {viol}", 10), ], body_size=10)
        y += Inches(0.705)
    card(s, Inches(8.46), Inches(1.70), Inches(4.26), Inches(1.86), "neut",
         title="五项必报输出", lines=["① 仪器覆盖率 O_j / N_j",
                                      "② 观测到的违规 X_j",
                                      "③ 不确定性（整群感知区间）",
                                      "④ 执法成本（计入 E）",
                                      "⑤ 运行域 Ω 的排除项"])
    card(s, Inches(8.46), Inches(3.72), Inches(4.26), Inches(2.14), "warn",
         title="三阈值与保守方向", lines=[
             "覆盖率下限 γ_j、违规率上限 ν_j、",
             "监测灵敏度下限 η_j，均在收集前声明。",
             "**未覆盖事件一律按违规计**：",
             "最坏率 (X_j + M_j)/N_j 的单侧 95% 上界须低于 ν_j。",
             "盲化诊断注入校准灵敏度；",
             "自报的 Harness trace 不能自证完备。"])
    footer(s, 11, "论文 §10 · 通过六项阈值不证明不存在第七条耦合渠道，只证明预登记的威胁模型被测得足够好")
    return s


def s12(prs):
    return figure_slide(
        prs, 12, 6, "Cluster-Period 随机交叉试验", "全文唯一的随机化单位", "实验设计", "phys",
        [("phys", "实验单位", ["cluster-period：隔离的 worker 池 / 租户分区 /",
                               "部署单元，在一个 (c,s) 下运行固定周期。",
                               "周期内请求是重复观测，非独立随机化单位。"], 1.62),
         ("logic", "指派与平衡", ["在集群与失效模式区块内随机化；",
                                  "对抗平衡序列使每个处理早晚位置相当，",
                                  "以便估计差异性结转。"], 1.42),
         ("data", "重置哨兵（四项全过才进下一周期）", ["前一周期契约身份清零、清单哈希精确匹配、",
                                                       "哨兵 log-p95 与全新起点的区间受界、",
                                                       "哨兵全需求满足风险差受界。"], 1.52),
         ("warn", "E 是边际量", ["配对回放固定 model/tool 输出，故「准入拒绝改变",
                                 "后续轨迹」的机会成本不计入 E，而归入 Q。"], 1.30)],
        "论文 §11 · 完整序列主估计 + 分配级 ITT 敏感性 + tipping-point；负载导致的失败是结局，不是排除项")


def s13(prs):
    s = blank(prs)
    heading(s, "四态判决规则", "第四态的存在是因为六项义务是合取", tag="判决", tagpal="phys")
    states = [("data", "supported within Ω",
               ["六项义务全部通过覆盖率、灵敏度校准与最坏情况违规上界；",
                "90% 交互区间落在 [−m_R, m_R] 内；每个语义下界超过 −m_Q；每项开销上界低于预算分量。"]),
              ("warn", "falsified within Ω",
               ["合格性与校准足以判决，且交互区间整体越出 margin，或语义上界低于 −m_Q，",
                "或观测违规下界超过上限，或执法开销下界超出预算。"]),
              ("phys", "conditional-engineering result  ★",
               ["一项或多项义务违规率超阈 → 报告完整结局估计 + 越阈义务及其观测值，**对猜想不作任何方向的判决**。",
                "它刻画 runtime 的工程状态，而非 P1 的科学地位；可发表、有信息量，且其违规率为后续研究提供剂量轴。"]),
              ("neut", "inconclusive",
               ["仪器覆盖或监测校准不足、区间无法分辨 margin、重置后仍有污染、",
                "预声明合格性规则失败，或所需量无法以计划的整群不确定性估计。"])]
    y = Inches(1.70)
    for pal, name, body in states:
        card(s, Inches(0.62), y, Inches(12.1), Inches(1.16), pal, title=name,
             title_size=14, lines=body, body_size=10)
        y += Inches(1.27)
    card(s, Inches(0.62), Inches(6.30), Inches(12.1), Inches(0.62), "warn",
         lines=[("判决规则区分「主张失败」与「未能测量」。功效不足导致的不显著交互是 inconclusive，不是支持。", 10.5)])
    footer(s, 13, "论文 §5.4 · 第四态由 v25 从企业篇提升为全局规则，使单个难测渠道不再永久锁死研究线")
    return s


def s14(prs):
    s = blank(prs)
    heading(s, "v27 形式化附录：四条命题", "回答一个很窄的问题：查询尚未知晓时，写时固定的表示能保住什么",
            tag="v27 新增", tagpal="contract")
    props = [("contract", "① 解码器无关的错误下界",
              ["亏损 Δ(f) = H(Y | Z, Q)；条件 Fano 给出 Δ(f) ≤ H_b(P_e) + P_e·log₂(M−1)。",
               "故 Δ(f) > 0 时，任何**仅用 (Z, Q)** 的解码器都有一个由该表示与任务共同决定的正错误下界。",
               "更强的 agent 只能减少解码器的次优性，无法重建 (Z, Q) 中已不存在的区分。",
               "注意：不等式给的是**下界，不是观测错误的点估计**，等号也不必成立。"]),
             ("logic", "② 划分判据、最小码率与查询单调性",
              ["确定编码器诱导源字母表的划分 π_f；每个查询诱导任务划分 π_q，π★ 是全部 π_q 的共同细化。",
               "**充分 ⟺ π_f 的每个块都被包含在 π★ 的某个块内**（即 π_f 细化 π★）；由此 H(Z) ≥ H(π★)，",
               "且 π★ 的块标识达到等号，故 R_min = H(π★)。扩大准入查询族只会让 π★ 更细，R_min 不可能下降。",
               "推论：**比特数本身不是保真保证** —— 表示的位长不变，任务族扩张后仍可能变得不充分。"]),
             ("phys", "③ 线性—对数分离（在源回退可用的前提下）",
              ["坐标查询族：任何精确回答全部坐标查询的写时表示满足 H(Z) ≥ k bit；",
               "若原始版本化源在读时仍可寻址，则一次「发地址 + 返回一个证据位」为 ⌈log₂k⌉ + 1 bit，比值 ≥ Θ(k / log₂k)。",
               "**是 k 的线性—对数关系，不是 k 的指数关系** —— 只有把地址长度 n = log₂k 当作独立参数时才呈指数。",
               "比较的是两种不同资源：未把原始存储、索引维护、授权检查与访问时延记作零成本。"]),
             ("data", "④ 接口不可区分性（统一多种失效模式）",
              ["任何存储/检索接口在源状态上诱导一个等价关系；若等价的两个状态对某准入查询答案不同，",
               "则受限于该接口的下游解码器不可能精确充分。粗检索单元、顺序丢失、超预算切片都是它的实例：",
               "顺序丢失是同一现象的群作用版本 —— n 个不同项在未知均匀排序下，置换不变表示丢弃 log₂(n!) bit 顺序身份（仅当顺序与任务相关时才算不足）。",
               "**可寻址性是整个接口与回退路径的性质，不是贴在某个算法上的标签** —— 保留稳定源标识并支持精确回退的 embedding 系统就没有该缺陷。"])]
    y = Inches(1.62)
    for pal, t, body in props:
        card(s, Inches(0.62), y, Inches(12.1), Inches(1.27), pal, title=t,
             title_size=12.5, lines=body, body_size=9.3)
        y += Inches(1.34)
    footer(s, 14, "论文附录 §A.1–A.4 · 随机编码、查询条件化物化、近似答案关系与连续源都需要不同的充分性定义，均在这些命题范围之外")
    return s


def s15(prs):
    s = blank(prs)
    heading(s, "形式化附录导出的五条设计约束", "以及一条明确的证伪边界", tag="v27 新增", tagpal="contract")
    cons = [("① 每个衍生物都必须保留可解析的来源与快照身份", "logic"),
            ("② 系统必须声明某个表示或策略被主张充分的**任务族**", "logic"),
            ("③ 索引与摘要应按其保留的**区分与证据路径**评价，而不是只看压缩比", "contract"),
            ("④ 回退不只需要留下字节 —— 还需可执行访问、授权、新鲜度、来源与版本兼容", "phys"),
            ("⑤ 任务族扩大或源分布改变即触发重新验证，因为最小充分划分可能已变", "data")]
    y = Inches(1.70)
    for t, pal in cons:
        card(s, Inches(0.62), y, Inches(12.1), Inches(0.62), pal,
             lines=[(t, 11)], body_size=11)
        y += Inches(0.70)
    card(s, Inches(0.62), Inches(5.28), Inches(12.1), Inches(1.14), "warn",
         title="附录明确不主张的四件事（必须一起讲）", lines=[
             "① 不证明 Data Wiki 降低时延或成本；② 不证明其索引优于直接读取；",
             "③ **不证明 embedding 检索普遍更差**；④ 不证明源回退在运维上可靠。",
             "这四件事属于 P16 的 held-out 对照与全生命周期成本测量，不属于证明。"])
    card(s, Inches(0.62), Inches(6.54), Inches(12.1), Inches(0.62), "neut",
         lines=[("证明只建立一条边界：在所述有限、确定性假设下，一旦接口合并了某个准入任务必须区分的两个源状态，"
                 "下游智能本身无法恢复精确性。", 10.5)])
    footer(s, 15, "论文附录 §A.5")
    return s


def s16(prs):
    s = blank(prs)
    heading(s, "novelty 边界", "最该主动承认的一篇，就在本文前提上", tag="划界", tagpal="warn")
    card(s, Inches(0.62), Inches(1.70), Inches(12.1), Inches(1.52), "warn",
         title="arXiv 2605.26112 — Scaling the Harness in Agentic AI（Shangding Gu, UC Berkeley）",
         title_size=13.5, lines=[
             "该文主张下一个瓶颈是 system scaling 而非 model scaling，并把 foundation model 周围的结构化执行层",
             "称为应被当作 “a first-class object of design, evaluation, and optimization” 的对象 —— 这就是本文所称的 Harness。",
             "→ **把 harness 认定为可扩展性的关键层，不是本文的贡献。** 本文剩余空间是承诺它的一条可证伪推论："
             "Ω 内 capability/capacity 可分离且执法成本有界，并给出随机化设计、margin、义务与四态判决。"])
    rows = [("contract", "vs Falsifiable Release Gates（2607.13070）",
             "其不变量治理 action-to-effector 路径，本文治理 capability-to-capacity 绑定；其 0.021 ms/请求被引为**下界**（只含不变量求值）。"),
            ("logic", "vs Five-Plane（2606.12320）",
             "治理视角的横切分面，可作本架构安全层，不建立可分离性主张。"),
            ("phys", "vs Skillware / SkillCorpus / 跨层错配 / Agent Skill Security",
             "支撑「Skill 是需行为性验证的受治资产」，不证明任何具体 Skill 安全或可移植。"),
            ("data", "vs GRACE（2607.09175）",
             "结构化使验证局部化；IR 把该原理用到它未涉及的边界 —— 源侧与产出侧 registry 的耦合。"),
            ("neut", "vs DeltaBox（2605.22781）",
             "毫秒级 sandbox checkpoint/rollback 是 reset 哨兵可能达成的**使能机制**，不是哨兵可达成的证据。")]
    y = Inches(3.38)
    for pal, t, body in rows:
        card(s, Inches(0.62), y, Inches(12.1), Inches(0.66), pal,
             lines=[(t, 10.5), (body, 9.4)], body_size=9.4)
        y += Inches(0.72)
    footer(s, 16, "论文 §4 · 2605.26112 长期未被识别，直接原因是本地 PDF 文件名标题写错（详见 papers/README.md 的核验说明）")
    return s


def s17(prs):
    s = blank(prs)
    heading(s, "证据边界：本文没有实测结果", "这一页不能从报告里删掉", tag="必读", tagpal="warn")
    card(s, Inches(0.62), Inches(1.72), Inches(12.1), Inches(1.16), "warn",
         title="本文不提供什么", title_size=14, lines=[
             "没有完成的 runtime 实现、没有实测的执法预算、没有数据集、没有 benchmark 结果、",
             "没有关于某个部署可支撑规模的主张。图是责任与协议视图，不是截图或实现证据。"])
    rows = [("supported 不等于「证明独立」",
             "只表示在预登记的 Ω、margin、功效与仪器条件内未触发任何证伪标准。任何更强外推都不成立。"),
            ("六项义务合取可能贵到不可测",
             "严格的覆盖规则降低假支持，但会把有用的部署推向 inconclusive 或 conditional-engineering。外部效度因此存疑。"),
            ("组织与治理侧最缺证据",
             "四类责任对象的 owner 划分、争议升级与跨企业可迁移性，尚未经任何组织研究检验 —— 这是全文最弱的一环。"),
            ("与现有编排框架的逐项对比仍未落地",
             "自 v20 评审提出至今未做；论文已在 §13.4 把它登记为评估程序中最显眼的缺口。")]
    y = Inches(3.10)
    for t, body in rows:
        card(s, Inches(0.62), y, Inches(12.1), Inches(0.80), "neut", title=t,
             title_size=12, lines=[(body, 9.6)], body_size=9.6)
        y += Inches(0.90)
    footer(s, 17, "论文 §14 · 证据基础是分层的：经典基础是稳定锚点，同期文献只支撑动机，不验证本架构")
    return s


def s18(prs):
    s = blank(prs)
    heading(s, "下一步：最小可行研究", "目的不是检验 P1，而是先验证协议可执行", tag="行动", tagpal="data")
    card(s, Inches(0.62), Inches(1.72), Inches(5.92), Inches(2.60), "data",
         title="单集群 2×2 可行性 pilot", title_size=14, lines=[
             "刻意声明一个很小的 Ω：单租户、一个 workload family、",
             "一个模型与 tokenizer 版本、一份策略快照、",
             "恰好两个兼容 Scaffold class（不做拓扑扫描）。",
             "能力轴用为实验撰写的合成 Skill，使声明面的 ground truth 可得。",
             "复制次数花在每格的 cluster-period 上，而不是更多格子。"])
    card(s, Inches(6.80), Inches(1.72), Inches(5.92), Inches(2.60), "warn",
         title="不能省的是仪器", title_size=14, lines=[
             "六项义务全部仪器化、六个违规率全部报告 ——",
             "即使结局对比功效不足，这些率本身就是主产物。",
             "**不要求六项全过**：完整中介率非零而其余五项为零，",
             "会得到一个 conditional-engineering 结果，",
             "它准确指出真实 runtime 留下了哪条渠道、关闭它要多少代价。"])
    card(s, Inches(0.62), Inches(4.50), Inches(12.1), Inches(1.06), "neut",
         title="预期产出（刻意有限且具体）", lines=[
             "六个违规率及其仪器、跨两个 Scaffold class 的语义区间、两点容量响应、",
             "带 achieved power 的交互估计，以及关闭每个越阈义务所需成本的账。"])
    card(s, Inches(0.62), Inches(5.72), Inches(12.1), Inches(1.10), "contract",
         title="这个规模能与不能回答什么", lines=[
             "**不能**建立可分离性，也不应声称建立了。",
             "**能**回答「六项义务在一个运行中的系统里究竟是否可仪器化」—— 这一点目前未知，且是所有更大规模研究的前提。"])
    footer(s, 18, "论文 §11.5")
    return s


def s19(prs):
    s = blank(prs)
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.22), SH)
    band.fill.solid()
    band.fill.fore_color.rgb = C("contract", 1)
    band.line.fill.background()
    band.shadow.inherit = False
    textbox(s, Inches(0.95), Inches(2.60), Inches(11.2), Inches(1.0),
            [("讨论", 38, True, INK)])
    textbox(s, Inches(0.95), Inches(3.70), Inches(11.4), Inches(2.0),
            [("论文    academy/agentic-runtime-preprint/output/pdf/…_v27.pdf（37 页 / 8 图 / 45 引用）", 12, False, MUT),
             ("正文    academy/agentic-runtime-preprint/paper_source/main.tex", 12, False, MUT),
             ("变更史  academy/agentic-runtime-preprint/CHANGELOG.md（v25 合回单篇 → v26 证据审计 → v27 形式化附录）", 12, False, MUT),
             ("大纲    academy/v25-outline.md", 12, False, MUT),
             ("评审 deck  academy/architecture-review-v25-zh.pptx", 12, False, MUT),
             ("本 deck    academy/paper-report-v27-zh.pptx（构建脚本 paper_report_build.py）", 12, False, MUT),
             ("笔记    Obsidian / Agentic-Runtime-参考架构（唯一权威，repo notes/ 为镜像）", 12, False, MUT)])
    return s


def main():
    prs = Presentation()
    prs.slide_width, prs.slide_height = SW, SH
    for fn in (s01, s02, s03, s04, s05, s06, s07, s08, s09,
               s10, s11, s12, s13, s14, s15, s16, s17, s18, s19):
        fn(prs)
    out = os.path.join(HERE, "paper-report-v27-zh.pptx")
    prs.save(out)
    print(f"wrote {out}  ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
